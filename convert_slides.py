import os
import re
import html
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.xmlchemy import OxmlElement

class HTMLToTextParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.text_lines = []
        self.current_line = []
        self.ignore_tags = {'svg', 'path', 'g', 'defs', 'marker', 'style', 'script'}
        self.ignore_depth = 0

    def handle_starttag(self, tag, attrs):
        if tag in self.ignore_tags:
            self.ignore_depth += 1
        elif tag in ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'div', 'li', 'tr', 'br']:
            self.flush_line()

    def handle_endtag(self, tag):
        if tag in self.ignore_tags:
            self.ignore_depth -= 1
        elif tag in ['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'div', 'li', 'tr', 'br']:
            self.flush_line()

    def handle_data(self, data):
        if self.ignore_depth == 0:
            text = data.strip()
            if text:
                self.current_line.append(text)

    def flush_line(self):
        if self.current_line:
            self.text_lines.append(" ".join(self.current_line))
            self.current_line = []

    def get_text_lines(self):
        self.flush_line()
        return [line for line in self.text_lines if line.strip()]

def parse_slide_content(slide_raw):
    lines = slide_raw.split('\n')
    
    notes_lines = []
    body_raw_lines = []
    in_note_block = False
    
    note_regex_html_start = re.compile(r'<!--\s*Note:', re.IGNORECASE)
    note_regex_html_end = re.compile(r'-->')
    note_regex_line = re.compile(r'^Note:\s*(.*)', re.IGNORECASE)
    
    for line in lines:
        trimmed = line.strip()
        
        # HTML Comment Note
        if note_regex_html_start.search(trimmed):
            in_note_block = True
            content = re.sub(r'<!--\s*Note:\s*', '', trimmed, flags=re.IGNORECASE)
            if '-->' in content:
                content = re.sub(r'\s*-->$', '', content)
                in_note_block = False
            if content.strip():
                notes_lines.append(content.strip())
            continue
        elif in_note_block:
            if '-->' in trimmed:
                content = re.sub(r'\s*-->$', '', trimmed)
                in_note_block = False
                if content.strip():
                    notes_lines.append(content.strip())
            else:
                notes_lines.append(line)
            continue
            
        # Single-line "Note:" text
        note_match = note_regex_line.match(trimmed)
        if note_match:
            notes_lines.append(note_match.group(1).strip())
            continue
            
        body_raw_lines.append(line)
        
    notes_text = '\n'.join(notes_lines).strip()
    body_text = '\n'.join(body_raw_lines).strip()
    
    # Check if the slide text has custom HTML elements
    if '<div' in body_text or '<span' in body_text or '<table' in body_text:
        parser = HTMLToTextParser()
        parser.feed(body_text)
        parsed_lines = parser.get_text_lines()
    else:
        parsed_lines = body_raw_lines

    title = None
    body_elements = []
    first_header_found = False
    
    for line in parsed_lines:
        trimmed = line.strip()
        if not trimmed:
            continue
            
        # Check standard headers
        header_match = re.match(r'^(#{1,6})\s+(.*)', trimmed)
        if header_match:
            header_level = len(header_match.group(1))
            header_content = header_match.group(2).strip()
            # Clean markdown formatting from headers
            header_content = re.sub(r'[\*\_\`]', '', header_content)
            if not first_header_found:
                title = header_content
                first_header_found = True
            else:
                body_elements.append({
                    'type': 'header',
                    'level': header_level,
                    'text': header_content
                })
            continue
            
        # Check if line already starts with star or bullet (extracted from HTML/CSS)
        if trimmed.startswith('✦') or trimmed.startswith('•'):
            char = trimmed[0]
            content = trimmed[1:].strip()
            body_elements.append({
                'type': 'bullet',
                'level': 0 if char == '✦' else 1,
                'text': content
            })
            continue
            
        # Check bullet points in Markdown
        list_match = re.match(r'^(\s*)([\*\-\+])\s+(.*)', line)
        if list_match:
            spaces = list_match.group(1)
            content = list_match.group(3).strip()
            indent_level = len(spaces) // 2
            body_elements.append({
                'type': 'bullet',
                'level': indent_level,
                'text': content
            })
            continue
            
        # Check numbered points in Markdown
        num_list_match = re.match(r'^(\s*)(\d+)\.\s+(.*)', line)
        if num_list_match:
            spaces = num_list_match.group(1)
            num = num_list_match.group(2)
            content = num_list_match.group(3).strip()
            indent_level = len(spaces) // 2
            body_elements.append({
                'type': 'numbered',
                'number': num,
                'level': indent_level,
                'text': content
            })
            continue
            
        # Normal text paragraph
        body_elements.append({
            'type': 'paragraph',
            'level': 0,
            'text': trimmed
        })
        
    return title, body_elements, notes_text

def add_rich_text(paragraph, text, default_color, highlight_color, font_name='Microsoft JhengHei'):
    # Highlight markdown bold text **bold**
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            content = part[2:-2]
            run = paragraph.add_run()
            run.text = content
            run.font.bold = True
            run.font.color.rgb = highlight_color
            run.font.name = font_name
        else:
            if part:
                run = paragraph.add_run()
                run.text = part
                run.font.color.rgb = default_color
                run.font.name = font_name

def apply_slide_theme(prs, slide, is_spiritual=False):
    background = slide.background
    fill = background.fill
    fill.solid()
    if is_spiritual:
        fill.fore_color.rgb = RGBColor(253, 252, 249) # Off-white
        border_color = RGBColor(14, 165, 233) # Sky-blue border
    else:
        fill.fore_color.rgb = RGBColor(11, 15, 25) # Dark Navy
        border_color = RGBColor(14, 165, 233) # Sky-blue border
    
    # Rounded slide border shape (representing .slide-border from original style)
    border_margin = Inches(0.1)
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        border_margin, border_margin,
        prs.slide_width - 2 * border_margin,
        prs.slide_height - 2 * border_margin
    )
    border.fill.background()
    border.line.color.rgb = border_color
    border.line.width = Pt(0.5)
    if border.adjustments:
        border.adjustments[0] = 0.02

def add_title_underline(slide, is_spiritual=False):
    # Underline representing bottom-border on h2 titles
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5), Inches(11.733), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-blue line
    line.line.fill.background()

def add_background_decorations(prs, slide, is_spiritual=False):
    if not is_spiritual:
        return
        
    import random
    icons = ['✦', '☥', '♥', '✦', '☥', '♥']
    num_icons = random.randint(6, 9)
    sw = prs.slide_width
    sh = prs.slide_height
    faded_gold = RGBColor(244, 240, 230) # Extremely faded gold on off-white background
    
    for _ in range(num_icons):
        icon = random.choice(icons)
        margin_x = Inches(1.5)
        margin_y = Inches(1.0)
        
        # Position near the margins
        if random.random() < 0.6:
            x = random.uniform(Inches(0.2), margin_x) if random.random() < 0.5 else random.uniform(sw - margin_x, sw - Inches(0.2))
            y = random.uniform(Inches(0.2), sh - Inches(0.2))
        else:
            x = random.uniform(Inches(0.2), sw - Inches(0.2))
            y = random.uniform(Inches(0.2), margin_y) if random.random() < 0.5 else random.uniform(sh - margin_y, sh - Inches(0.2))
            
        font_size = random.randint(24, 48)
        
        tb = slide.shapes.add_textbox(x, y, Inches(1.0), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = icon
        run.font.name = 'M PLUS Rounded 1c'
        run.font.size = Pt(font_size)
        run.font.color.rgb = faded_gold
        
        tb.rotation = random.randint(0, 360)

def draw_hub_and_spoke_diagram(prs, slide):
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "全球科技巨頭到底在爭奪什麼？"
    run_title.font.name = 'Microsoft JhengHei'
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(192, 132, 252)
    
    add_title_underline(slide)
    
    # Coordinates for center hub
    cx = Inches(6.666)
    cy = Inches(4.3)
    hub_radius = Inches(1.0)
    
    # 1. Draw Hub (Center Circle)
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - hub_radius, cy - hub_radius,
        hub_radius * 2, hub_radius * 2
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = RGBColor(79, 70, 229) # Indigo-600
    hub.line.color.rgb = RGBColor(241, 245, 249)
    hub.line.width = Pt(2)
    
    # Text inside Hub
    tf_hub = hub.text_frame
    tf_hub.word_wrap = True
    p_hub = tf_hub.paragraphs[0]
    p_hub.alignment = PP_ALIGN.CENTER
    run_hub = p_hub.add_run()
    run_hub.text = "認知入口"
    run_hub.font.name = 'Microsoft JhengHei'
    run_hub.font.size = Pt(18)
    run_hub.font.bold = True
    run_hub.font.color.rgb = RGBColor(255, 255, 255)
    
    p_hub_sub = tf_hub.add_paragraph()
    p_hub_sub.alignment = PP_ALIGN.CENTER
    run_hub_sub = p_hub_sub.add_run()
    run_hub_sub.text = "Cognitive Portal"
    run_hub_sub.font.name = 'Microsoft JhengHei'
    run_hub_sub.font.size = Pt(10)
    run_hub_sub.font.color.rgb = RGBColor(224, 231, 255)
    
    # 2. Outer Spokes Data
    spokes = [
        ("真靈光", "身心靈一站式認知", RGBColor(251, 191, 36), cx - Inches(1.5), cy - Inches(2.7), cx, cy - Inches(2.7) + Inches(1.0), cx, cy - hub_radius),
        ("OpenAI", "未來知識入口", RGBColor(129, 140, 248), cx - Inches(5.2), cy - Inches(0.5), cx - Inches(5.2) + Inches(3.0), cy, cx - hub_radius, cy),
        ("Google", "資訊入口", RGBColor(52, 211, 153), cx + Inches(2.2), cy - Inches(0.5), cx + Inches(2.2), cy, cx + hub_radius, cy),
        ("Apple", "用戶體驗入口", RGBColor(239, 68, 68), cx - Inches(4.3), cy + Inches(1.5), cx - Inches(4.3) + Inches(2.5), cy + Inches(1.5), cx - Inches(0.7), cy + Inches(0.7)),
        ("Meta", "社交認知入口", RGBColor(244, 114, 182), cx + Inches(1.3), cy + Inches(1.5), cx + Inches(1.3) + Inches(0.5), cy + Inches(1.5), cx + Inches(0.7), cy + Inches(0.7))
    ]
    
    for name, sub, color, bx, by, lx1, ly1, lx2, ly2 in spokes:
        # Draw Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bx, by, Inches(3.0), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate-900
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        if box.adjustments:
            box.adjustments[0] = 0.05
            
        tf_box = box.text_frame
        tf_box.word_wrap = True
        p_name = tf_box.paragraphs[0]
        p_name.alignment = PP_ALIGN.CENTER
        run_name = p_name.add_run()
        run_name.text = name
        run_name.font.name = 'Microsoft JhengHei'
        run_name.font.bold = True
        run_name.font.size = Pt(16)
        run_name.font.color.rgb = RGBColor(255, 255, 255)
        
        p_sub = tf_box.add_paragraph()
        p_sub.alignment = PP_ALIGN.CENTER
        run_sub = p_sub.add_run()
        run_sub.text = sub
        run_sub.font.name = 'Microsoft JhengHei'
        run_sub.font.size = Pt(11)
        run_sub.font.color.rgb = color
        
        # Draw Arrow line (spoke to hub)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, lx1, ly1, lx2, ly2)
        conn.line.color.rgb = color
        conn.line.width = Pt(1.5)
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        
        # Add XML arrowhead
        ln = conn._element.spPr.get_or_add_ln()
        endArrow = OxmlElement('a:endArrow')
        endArrow.set('type', 'triangle')
        endArrow.set('w', 'med')
        endArrow.set('len', 'med')
        ln.append(endArrow)

def draw_sally_integration_diagram(prs, slide):
    # Add slide title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    run_title = p_title.add_run()
    run_title.text = "整合型定位：專長 + 身心靈 = 質變方程式"
    run_title.font.name = 'Microsoft JhengHei'
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(192, 132, 252)
    
    add_title_underline(slide)
    
    # 1. Left Box
    box_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.8), Inches(4.5), Inches(1.8))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = RGBColor(30, 27, 75)
    box_left.line.color.rgb = RGBColor(14, 165, 233) # Sky-blue Left Box Line
    box_left.line.width = Pt(2)
    if box_left.adjustments:
         box_left.adjustments[0] = 0.05
         
    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    p_l1 = tf_l.paragraphs[0]
    p_l1.alignment = PP_ALIGN.CENTER
    p_l1.space_after = Pt(8)
    run_l1 = p_l1.add_run()
    run_l1.text = "您原有的專業"
    run_l1.font.name = 'Microsoft JhengHei'
    run_l1.font.bold = True
    run_l1.font.size = Pt(18)
    run_l1.font.color.rgb = RGBColor(255, 255, 255)
    
    p_l2 = tf_l.add_paragraph()
    p_l2.alignment = PP_ALIGN.CENTER
    run_l2 = p_l2.add_run()
    run_l2.text = "如：推拿整脊、美容芳療、保險經紀、幼教護理等實體專長"
    run_l2.font.name = 'Microsoft JhengHei'
    run_l2.font.size = Pt(13)
    run_l2.font.color.rgb = RGBColor(199, 210, 254)
    
    # 2. Plus Sign
    plus_box = slide.shapes.add_textbox(Inches(5.9), Inches(2.35), Inches(0.5), Inches(0.5))
    tf_p = plus_box.text_frame
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.CENTER
    run_p = p_p.add_run()
    run_p.text = "+"
    run_p.font.name = 'Microsoft JhengHei'
    run_p.font.bold = True
    run_p.font.size = Pt(36)
    run_p.font.color.rgb = RGBColor(14, 165, 233) # Sky-blue Plus Sign
    
    # 3. Right Box
    box_right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.8), Inches(4.5), Inches(1.8))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = RGBColor(59, 7, 100)
    box_right.line.color.rgb = RGBColor(14, 165, 233) # Sky-blue Right Box Line
    box_right.line.width = Pt(2)
    if box_right.adjustments:
         box_right.adjustments[0] = 0.05
         
    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    p_r1 = tf_r.paragraphs[0]
    p_r1.alignment = PP_ALIGN.CENTER
    p_r1.space_after = Pt(8)
    run_r1 = p_r1.add_run()
    run_r1.text = "熱愛的身心靈能力"
    run_r1.font.name = 'Microsoft JhengHei'
    run_r1.font.bold = True
    run_r1.font.size = Pt(18)
    run_r1.font.color.rgb = RGBColor(255, 255, 255)
    
    p_r2 = tf_r.add_paragraph()
    p_r2.alignment = PP_ALIGN.CENTER
    run_r2 = p_r2.add_run()
    run_r2.text = "如：靈氣、生命數字、情緒溝通、奇門占卜、正念引導等"
    run_r2.font.name = 'Microsoft JhengHei'
    run_r2.font.size = Pt(13)
    run_r2.font.color.rgb = RGBColor(233, 213, 255)

    # 4. Down Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(5.95), Inches(3.9), Inches(0.4), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-blue Arrow
    arrow.line.fill.background()
    
    # 5. Bottom Result Box
    box_bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.166), Inches(4.7), Inches(8.0), Inches(1.8))
    box_bottom.fill.solid()
    box_bottom.fill.fore_color.rgb = RGBColor(15, 23, 42)
    box_bottom.line.color.rgb = RGBColor(14, 165, 233) # Sky-blue Bottom Box Line
    box_bottom.line.width = Pt(2.5)
    if box_bottom.adjustments:
         box_bottom.adjustments[0] = 0.05
         
    tf_b = box_bottom.text_frame
    tf_b.word_wrap = True
    p_b1 = tf_b.paragraphs[0]
    p_b1.alignment = PP_ALIGN.CENTER
    p_b1.space_after = Pt(8)
    run_b1 = p_b1.add_run()
    run_b1.text = "✨ 整合型定位 (質變與差異化) ✨"
    run_b1.font.name = 'Microsoft JhengHei'
    run_b1.font.bold = True
    run_b1.font.size = Pt(20)
    run_b1.font.color.rgb = RGBColor(14, 165, 233) # Sky-blue Text Color
    
    p_b2 = tf_b.add_paragraph()
    p_b2.alignment = PP_ALIGN.CENTER
    run_b2 = p_b2.add_run()
    run_b2.text = "拉高生命維度，在原有紅海競爭中脫穎而出，成為同業中唯一的發光點。"
    run_b2.font.name = 'Microsoft JhengHei'
    run_b2.font.size = Pt(14)
    run_b2.font.color.rgb = RGBColor(241, 245, 249)

def convert_html_to_pptx(html_path, pptx_path):
    print(f"Parsing {html_path}...")
    with open(html_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    # Extract textarea contents
    match = re.search(r'<textarea\s+id="editor"[^>]*>([\s\S]*?)<\/textarea>', content)
    if not match:
        print(f"Error: Could not find editor textarea in {html_path}")
        return
        
    markdown_content = html.unescape(match.group(1))
    slides_raw = markdown_content.split('\n---\n')
    
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    is_spiritual = ("身心靈" in html_path) or ("spiritual" in html_path)
    
    # Set theme colors & fonts
    if is_spiritual:
        text_color_header = RGBColor(184, 151, 90) # Gold
        text_color_subtitle = RGBColor(14, 165, 233) # Sky-blue Soft Accent
        text_color_body = RGBColor(60, 60, 60) # Dark Gray
        text_color_highlight = RGBColor(3, 105, 161) # Sky-blue Dark bold
        font_name = 'M PLUS Rounded 1c'
    else:
        text_color_header = RGBColor(192, 132, 252) # Light Purple / Violet
        text_color_subtitle = RGBColor(14, 165, 233) # Sky-blue Accent
        text_color_body = RGBColor(241, 245, 249) # Off-white
        text_color_highlight = RGBColor(56, 189, 248) # Sky-blue Light highlight
        font_name = 'Microsoft JhengHei'

    parsed_slides = []
    for idx, slide_raw in enumerate(slides_raw):
        title, body_elements, notes = parse_slide_content(slide_raw)
        
        is_cover = False
        if idx == 0:
            is_cover = True
        elif not any(el['type'] in ['bullet', 'numbered'] for el in body_elements) and len(body_elements) <= 2:
            is_cover = True
            
        parsed_slides.append({
            'title': title,
            'body_elements': body_elements,
            'notes': notes,
            'is_cover': is_cover,
            'raw': slide_raw
        })

    # Apply slide splitting logic
    final_slides = []
    for slide_data in parsed_slides:
        if slide_data['is_cover'] or not slide_data['body_elements']:
            final_slides.append(slide_data)
            continue
            
        body_elements = slide_data['body_elements']
        chunks = []
        current_chunk = []
        current_lines = 0
        
        for el in body_elements:
            text = el['text']
            clean_text = text.replace('**', '')
            char_limit = 43 if el['level'] > 0 else 38
            el_lines = max(1, (len(clean_text) + char_limit - 1) // char_limit)
            
            if current_lines + el_lines > 10 and current_chunk:
                chunks.append(current_chunk)
                current_chunk = [el]
                current_lines = el_lines
            else:
                current_chunk.append(el)
                current_lines += el_lines
                
        if current_chunk:
            chunks.append(current_chunk)
            
        total_pages = len(chunks)
        if total_pages <= 1:
            final_slides.append(slide_data)
            continue
            
        for idx, chunk in enumerate(chunks):
            page_num = idx + 1
            original_title = slide_data['title']
            new_title = f"{original_title} ({page_num}/{total_pages})" if original_title else f"({page_num}/{total_pages})"
            
            final_slides.append({
                'title': new_title,
                'body_elements': chunk,
                'notes': slide_data['notes'],
                'is_cover': slide_data['is_cover'],
                'raw': slide_data['raw']
            })

    for idx, slide_data in enumerate(final_slides):
        title = slide_data['title']
        body_elements = slide_data['body_elements']
        notes = slide_data['notes']
        is_cover = slide_data['is_cover']
        
        slide = prs.slides.add_slide(blank_layout)
        apply_slide_theme(prs, slide, is_spiritual)
        add_background_decorations(prs, slide, is_spiritual)
        
        # Put Speaker Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
            
        # Draw custom diagrams if title matches
        is_tech_giants = (title == "全球科技巨頭到底在爭奪什麼？") or ("全球科技巨頭到底在爭奪什麼" in slide_data['raw'])
        if is_tech_giants:
            draw_hub_and_spoke_diagram(prs, slide)
            continue
        elif title == "整合型定位：專長 + 身心靈 = 質變方程式":
            draw_sally_integration_diagram(prs, slide)
            continue
            
        if is_cover:
            # Centered layout for cover slides
            title_text = title if title else "投影片"
            subtitle_text = ""
            if body_elements:
                subtitle_text = "\n".join([el['text'] for el in body_elements])
                
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(14)
            run = p.add_run()
            run.text = title_text
            run.font.name = font_name
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = text_color_header
            
            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                run2 = p2.add_run()
                run2.text = subtitle_text
                run2.font.name = font_name
                run2.font.size = Pt(24)
                run2.font.color.rgb = text_color_subtitle
        else:
            # Content layout
            # Add Title
            title_text = title if title else ""
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(0.8))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
            
            p_title = tf_title.paragraphs[0]
            run_title = p_title.add_run()
            run_title.text = title_text
            run_title.font.name = font_name
            run_title.font.size = Pt(32)
            run_title.font.bold = True
            run_title.font.color.rgb = text_color_header
            
            add_title_underline(slide, is_spiritual)
            
            # Add Content Box
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            tf_body.margin_left = tf_body.margin_top = tf_body.margin_right = tf_body.margin_bottom = 0
            
            for b_idx, el in enumerate(body_elements):
                p = tf_body.paragraphs[0] if b_idx == 0 else tf_body.add_paragraph()
                p.level = min(el['level'], 4)
                
                # Determine font sizing, prefix and styling
                font_size = Pt(20 - 2 * p.level)
                
                if el['type'] == 'header':
                    p.space_before = Pt(14)
                    p.space_after = Pt(8)
                    run = p.add_run()
                    run.text = el['text']
                    run.font.name = font_name
                    run.font.bold = True
                    run.font.size = font_size + Pt(2)
                    run.font.color.rgb = text_color_subtitle
                elif el['type'] == 'bullet':
                    p.space_before = Pt(4)
                    p.space_after = Pt(8)
                    p.left_indent = Inches(0.4 * (p.level + 1))
                    p.first_line_indent = Inches(-0.25)
                    
                    bullet_char = '✦ ' if p.level == 0 else '• '
                    run_bullet = p.add_run()
                    run_bullet.text = bullet_char
                    run_bullet.font.name = font_name
                    run_bullet.font.bold = True
                    run_bullet.font.size = font_size
                    run_bullet.font.color.rgb = text_color_subtitle
                    
                    add_rich_text(p, el['text'], text_color_body, text_color_highlight, font_name)
                elif el['type'] == 'numbered':
                    p.space_before = Pt(4)
                    p.space_after = Pt(8)
                    p.left_indent = Inches(0.4 * (p.level + 1))
                    p.first_line_indent = Inches(-0.25)
                    
                    num_prefix = f"{el['number']}. "
                    run_num = p.add_run()
                    run_num.text = num_prefix
                    run_num.font.name = font_name
                    run_num.font.bold = True
                    run_num.font.size = font_size
                    run_num.font.color.rgb = text_color_subtitle
                    
                    add_rich_text(p, el['text'], text_color_body, text_color_highlight, font_name)
                else: # paragraph
                    p.space_before = Pt(6)
                    p.space_after = Pt(10)
                    add_rich_text(p, el['text'], text_color_body, text_color_highlight, font_name)
                    
                # Apply font sizes to all runs in paragraph
                for r in p.runs:
                    if r.text not in ['✦ ', '• ', f"{el.get('number', '')}. "]:
                        r.font.size = font_size

    prs.save(pptx_path)
    print(f"Successfully saved to {pptx_path}\n")

if __name__ == '__main__':
    base_dir = r"d:\Git\true-soul-light"
    files = [
        '世豐老師脆課程.html',
        '世豐老師脆課程完整版.html',
        '艾莫老師簡報.html',
        '莎莉身心靈簡報.html',
        '不用版控的/陳仲豪總監/總監課程簡報.html',
        '不用版控的/世豐脆課程/世豐脆課程第一堂_身心靈風格.html'
    ]
    
    for file_name in files:
        html_file = os.path.join(base_dir, file_name)
        pptx_name = file_name.replace('.html', '.pptx')
        pptx_file = os.path.join(base_dir, pptx_name)
        
        if os.path.exists(html_file):
            convert_html_to_pptx(html_file, pptx_file)
        else:
            print(f"Error: {html_file} does not exist.")
