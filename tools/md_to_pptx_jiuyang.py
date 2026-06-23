import re
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_jiuyang_markdown(filepath):
    """
    Parses the specialized jiuyang markdown file into logical slides.
    """
    if not os.path.exists(filepath):
        print(f"Error: File not found: {filepath}")
        return []

    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Normalize line endings
    content = content.replace('\r\n', '\n')
    
    # Split by "---" to get raw slide blocks
    raw_blocks = content.split('\n---')
    
    slides = []
    
    for block in raw_blocks:
        block = block.strip()
        if not block:
            continue
            
        # Stop processing once we reach the appendix
        if "📎 附錄" in block or "附錄一" in block or "附錄二" in block or "附錄三" in block:
            break
            
        is_appendix = "📎 附錄" in block
        
        if is_appendix:
            lines = block.split('\n')
            sub_blocks = []
            current_sub_block = []
            
            for line in lines:
                stripped = line.strip()
                if stripped.startswith('##') and not stripped.startswith('###'):
                    if current_sub_block:
                        sub_blocks.append("\n".join(current_sub_block))
                    current_sub_block = [line]
                else:
                    current_sub_block.append(line)
            if current_sub_block:
                sub_blocks.append("\n".join(current_sub_block))
                
            for sb in sub_blocks:
                sb = sb.strip()
                if sb:
                    slide_data = parse_single_slide(sb)
                    if slide_data:
                        slides.append(slide_data)
        else:
            slide_data = parse_single_slide(block)
            if slide_data:
                slides.append(slide_data)
                
    return slides

def parse_single_slide(block_text):
    lines = block_text.split('\n')
    title = ""
    subtitle = ""
    notes_lines = []
    content_lines = []
    
    # Find the slide title (first line starting with # or ##)
    title_line_idx = -1
    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith('#') or stripped.startswith('##'):
            title_line_idx = i
            raw_title = stripped.lstrip('#').strip()
            # Remove 🎞️ Slide X： prefix
            title = re.sub(r'^(?:🎞️\s*)?Slide\s*\d+\s*[：:]\s*', '', raw_title)
            title = title.strip()
            break
            
    if title_line_idx == -1:
        # Fallback if no title line found
        non_empty = [l.strip() for l in lines if l.strip()]
        if not non_empty:
            return None
        title = "前言：行銷與變現分享"
        start_idx = 0
    else:
        start_idx = title_line_idx + 1
        
    # Check if there are explicit subsections for screen content and speaker notes
    has_subsections = False
    for line in lines[start_idx:]:
        stripped = line.strip()
        if "###" in stripped and ("畫面文字" in stripped or "現場講稿" in stripped):
            has_subsections = True
            break
            
    in_notes = False
    in_content = not has_subsections
    
    for line in lines[start_idx:]:
        stripped = line.strip()
        if not stripped:
            if in_notes:
                notes_lines.append("")
            elif in_content:
                content_lines.append("")
            continue
            
        if "###" in stripped and ("現場講稿" in stripped or "🎙️" in stripped):
            in_notes = True
            in_content = False
            continue
        elif "###" in stripped and "畫面文字" in stripped:
            in_content = True
            in_notes = False
            continue
        elif stripped.startswith('###') and has_subsections:
            # Skip metadata subheadings, but keep in content if we don't know it
            pass
            
        if in_notes:
            # Strip blockquote prefix
            note_line = stripped
            if note_line.startswith('>'):
                note_line = note_line[1:].strip()
            if note_line.startswith('「') and note_line.endswith('」'):
                note_line = note_line[1:-1]
            notes_lines.append(note_line)
        elif in_content:
            content_lines.append(line)
            
    # Clean up empty lines from content
    while content_lines and not content_lines[0].strip():
        content_lines.pop(0)
    while content_lines and not content_lines[-1].strip():
        content_lines.pop()
        
    notes = "\n".join(notes_lines).strip()
    
    # Check if we have H2 subtitle inside content lines
    sub_candidate = ""
    remaining_content_lines = []
    for line in content_lines:
        stripped = line.strip()
        if stripped.startswith('##') and not stripped.startswith('###'):
            sub_candidate = stripped.lstrip('#').strip()
        else:
            remaining_content_lines.append(line)
            
    if sub_candidate:
        subtitle = sub_candidate
        content_lines = remaining_content_lines
        
    is_cover = False
    if title == "行銷界九陽神功" or "封面頁" in title or title == "行銷界九陽神功：3 ╳ 3 大健康產業終極商業矩陣" or "壓軸金句" in title or "結尾 CTA" in title or "CTA" in title:
        is_cover = True
        
    items = []
    table_data = None
    
    # Parse Table if any
    table_lines = [l.strip() for l in content_lines if l.strip().startswith('|')]
    if table_lines:
        headers = []
        rows = []
        for tl in table_lines:
            cells = [c.strip() for c in tl.split('|')[1:-1]]
            if not cells:
                continue
            if all(re.match(r'^[-:\s]+$', c) for c in cells):
                continue
            if not headers:
                headers = cells
            else:
                rows.append(cells)
        if headers:
            table_data = {'headers': headers, 'rows': rows}
            
    # Parse Items (Bullets)
    for line in content_lines:
        stripped = line.strip()
        if not stripped or stripped.startswith('|'):
            continue
            
        leading_spaces = len(line) - len(line.lstrip())
        level = 1 if leading_spaces >= 2 else 0
        
        if stripped.startswith('- ') or stripped.startswith('* '):
            items.append({'type': 'bullet', 'text': stripped[2:], 'level': level})
        elif re.match(r'^\d+\.\s+', stripped):
            match = re.match(r'^(\d+)\.\s+(.*)', stripped)
            items.append({'type': 'num_list', 'text': match.group(2), 'num': match.group(1), 'level': level})
        elif stripped.startswith('>'):
            text = stripped.lstrip('>').strip()
            if text.startswith('「') and text.endswith('」'):
                text = text[1:-1]
            items.append({'type': 'quote', 'text': text, 'level': level})
        else:
            items.append({'type': 'text', 'text': stripped, 'level': level})
            
    # Parse Cover metadata
    if is_cover:
        content_text = "\n".join(content_lines)
        if '核心大字' in content_text:
            match = re.search(r'核心大字[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                title = match.group(1).strip().replace('**', '').replace('「', '').replace('」', '')
                title = "\n".join([l.strip() for l in title.split('\n') if l.strip()])
        elif '主標題' in content_text:
            match = re.search(r'主標題[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                title = match.group(1).strip().replace('**', '')
                
        if '底字' in content_text:
            match = re.search(r'底字[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                subtitle = match.group(1).strip().replace('**', '')
        elif '副標題' in content_text:
            match = re.search(r'副標題[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                subtitle = match.group(1).strip().replace('**', '')
                
        annotation = ""
        if '標註' in content_text:
            match = re.search(r'標註[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                annotation = match.group(1).strip().replace('**', '')
        elif '右側卡片' in content_text:
            match = re.search(r'右側卡片[：:]\s*(.*?)(?=\n\s*-\s*\*\*|$)', content_text, re.DOTALL)
            if match:
                annotation = match.group(1).strip().replace('**', '')
                
        if annotation:
            subtitle = f"{subtitle}\n({annotation})" if subtitle else annotation

    # Detect Master/Saint Quotes (Dual Column Layout)
    has_dual_quotes = False
    quotes_data = []
    
    # 1. Detect by emoji 👑
    master_items = [it for it in items if '👑' in it['text']]
    if len(master_items) >= 2:
        has_dual_quotes = True
        for mit in master_items:
            raw_text = mit['text'].replace('👑', '').strip()
            name_match = re.match(r'^\s*\*\*(.*?)\*\*\s*[：:]\s*(.*)', raw_text)
            if name_match:
                quotes_data.append({
                    'name': name_match.group(1).strip(),
                    'quote': name_match.group(2).strip()
                })
            else:
                quotes_data.append({
                    'name': '大師語錄',
                    'quote': raw_text
                })
                
    # 2. Detect by explicit left/right column structure (e.g. Slide 8)
    if not has_dual_quotes:
        left_text = ""
        right_text = ""
        left_name = "西方大師"
        right_name = "東方聖人"
        is_left = False
        is_right = False
        
        for it in items:
            t = it['text']
            if '左半部分' in t or '左半邊' in t:
                is_left = True
                is_right = False
                continue
            elif '右半部分' in t or '右半邊' in t:
                is_right = True
                is_left = False
                continue
                
            if is_left:
                if '🖼️' in t:
                    name_clean = t.replace('🖼️', '').strip()
                    name_clean = re.sub(r'\(.*?\)', '', name_clean)
                    name_clean = name_clean.replace('**', '').split('高清')[0].split('肖像')[0].strip()
                    if name_clean:
                        left_name = name_clean
                elif '名言' not in t and t.strip():
                    left_text += t + "\n"
            elif is_right:
                if '🖼️' in t:
                    name_clean = t.replace('🖼️', '').strip()
                    name_clean = re.sub(r'\(.*?\)', '', name_clean)
                    name_clean = name_clean.replace('**', '').split('水墨')[0].split('肖像')[0].strip()
                    if name_clean:
                        right_name = name_clean
                elif '名言' not in t and t.strip():
                    right_text += t + "\n"
                    
        if left_text or right_text:
            has_dual_quotes = True
            left_text = left_text.replace('*', '').replace('➔', '➔\n').strip()
            right_text = right_text.replace('*', '').replace('➔', '➔\n').strip()
            quotes_data = [
                {'name': left_name, 'quote': left_text},
                {'name': right_name, 'quote': right_text}
            ]
            
    return {
        'title': title,
        'subtitle': subtitle,
        'is_cover': is_cover,
        'items': items,
        'table_data': table_data,
        'has_dual_quotes': has_dual_quotes,
        'quotes_data': quotes_data,
        'notes': notes
    }

def split_slides(slides_data):
    """
    Splits long content slides into multiple slides.
    Includes logic to prevent level-0 headers from being stranded as the last line of a slide.
    Also splits dual quotes slides that contain general lists or concepts.
    """
    new_slides = []
    for slide in slides_data:
        # Handle dual quotes slide split
        if slide['has_dual_quotes']:
            if "雙聖人" in slide['title'] or "Slide 8" in slide['title']:
                new_slides.append(slide)
                continue
                
            non_quote_content = []
            case_study_item = None
            
            for item in slide['items']:
                t = item['text']
                if '實戰案例對撞' in t or '實戰案例' in t:
                    case_study_item = item
                elif '大師佐證' in t or '雙大師佐證' in t or '👑' in t or '左半部分' in t or '右半部分' in t or '左半邊' in t or '右半邊' in t or '肖像' in t or '名言' in t:
                    continue
                else:
                    is_quote_part = False
                    for q in slide['quotes_data']:
                        if q['name'] in t or (len(q['quote']) > 10 and q['quote'][:10] in t):
                            is_quote_part = True
                            break
                    if not is_quote_part:
                        non_quote_content.append(item)
            
            if len(non_quote_content) > 0:
                # Slide Part 1: Standard layout (Core concept + list items)
                slide_part1 = {
                    'title': f"{slide['title']} (1/2)" if slide['title'] else "(1/2)",
                    'subtitle': slide['subtitle'],
                    'is_cover': False,
                    'items': non_quote_content,
                    'table_data': None,
                    'has_dual_quotes': False,
                    'quotes_data': [],
                    'notes': slide['notes']
                }
                
                # Slide Part 2: Dual quotes layout (Left/Right cards + bottom case study Callout)
                part2_items = []
                if case_study_item:
                    part2_items.append(case_study_item)
                    
                slide_part2 = {
                    'title': f"{slide['title']} (2/2)" if slide['title'] else "(2/2)",
                    'subtitle': '',
                    'is_cover': False,
                    'items': part2_items,
                    'table_data': None,
                    'has_dual_quotes': True,
                    'quotes_data': slide['quotes_data'],
                    'notes': slide['notes']
                }
                
                new_slides.append(slide_part1)
                new_slides.append(slide_part2)
                continue
            else:
                new_slides.append(slide)
                continue
                
        # For non-dual-quotes content slides
        if slide['is_cover'] or slide['table_data']:
            new_slides.append(slide)
            continue
            
        items = slide['items']
        if not items:
            new_slides.append(slide)
            continue
            
        chunks = []
        current_chunk = []
        current_lines = 0
        max_lines_per_slide = 7
        
        for item in items:
            text = item['text']
            clean_text = text.replace('**', '')
            char_limit = 43 if item['level'] > 0 else 38
            item_lines = max(1, (len(clean_text) + char_limit - 1) // char_limit)
            
            if current_lines + item_lines > max_lines_per_slide and current_chunk:
                # Rule: Do not leave a level 0 header as the last line
                if current_chunk[-1]['level'] == 0:
                    last_header = current_chunk.pop()
                    chunks.append(current_chunk)
                    current_chunk = [last_header, item]
                    current_lines = 0
                    for x in current_chunk:
                        x_clean = x['text'].replace('**', '')
                        x_char_limit = 43 if x['level'] > 0 else 38
                        current_lines += max(1, (len(x_clean) + x_char_limit - 1) // x_char_limit)
                else:
                    chunks.append(current_chunk)
                    current_chunk = [item]
                    current_lines = item_lines
            else:
                current_chunk.append(item)
                current_lines += item_lines
                
        if current_chunk:
            chunks.append(current_chunk)
            
        total_pages = len(chunks)
        if total_pages <= 1:
            new_slides.append(slide)
            continue
            
        for idx, chunk in enumerate(chunks):
            page_num = idx + 1
            original_title = slide['title']
            new_title = f"{original_title} ({page_num}/{total_pages})" if original_title else f"({page_num}/{total_pages})"
            
            new_slides.append({
                'title': new_title,
                'subtitle': slide['subtitle'],
                'is_cover': slide['is_cover'],
                'items': chunk,
                'table_data': slide['table_data'],
                'has_dual_quotes': slide['has_dual_quotes'],
                'quotes_data': slide['quotes_data'],
                'notes': slide['notes']
            })
            
    return new_slides

def split_table_slides(slides_data):
    """
    Splits long table slides into multiple slides.
    """
    new_slides = []
    for slide in slides_data:
        if not slide['table_data']:
            new_slides.append(slide)
            continue
            
        headers = slide['table_data']['headers']
        rows = slide['table_data']['rows']
        
        max_table_rows = 6
        if len(rows) <= max_table_rows:
            new_slides.append(slide)
            continue
            
        chunks = [rows[i:i + max_table_rows] for i in range(0, len(rows), max_table_rows)]
        total_pages = len(chunks)
        
        for idx, chunk in enumerate(chunks):
            page_num = idx + 1
            original_title = slide['title']
            new_title = f"{original_title} ({page_num}/{total_pages})" if original_title else f"({page_num}/{total_pages})"
            
            new_slides.append({
                'title': new_title,
                'subtitle': slide['subtitle'],
                'is_cover': slide['is_cover'],
                'items': slide['items'],
                'table_data': {'headers': headers, 'rows': chunk},
                'has_dual_quotes': slide['has_dual_quotes'],
                'quotes_data': slide['quotes_data'],
                'notes': slide['notes']
            })
            
    return new_slides

def add_bold_text_run(paragraph, text, normal_color_rgb, bold_color_rgb):
    """
    Splits text by **bold** tags and adds formatted runs to paragraph.
    """
    parts = re.split(r'(\*\*.*?\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.color.rgb = bold_color_rgb
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.bold = False
            run.font.color.rgb = normal_color_rgb

def draw_callout_box(slide, text, top_inch, height_inch, bg_rgb, border_rgb, text_color_rgb):
    """
    Draws a styled callout box (rounded rectangle) on the slide.
    """
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), top_inch, Inches(11.73), height_inch
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_rgb
    box.line.color.rgb = border_rgb
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    
    add_bold_text_run(p, text, RGBColor(51, 65, 85), text_color_rgb)
    p.font.name = 'Microsoft JhengHei'
    p.font.size = Pt(16)

def create_presentation(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply premium light background (#FCFCFC)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(252, 252, 252)
        
        # Draw slide border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7)
        )
        border.fill.background()
        border.line.color.rgb = RGBColor(226, 232, 240) # Slate-200
        border.line.width = Pt(1.0)
        
        # Add speaker notes
        if slide_data['notes']:
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
            
        if slide_data['is_cover']:
            # Render Cover Slide
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 41, 59) # Slate-800
            p.alignment = PP_ALIGN.CENTER
            
            # Decorative centered line
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.03)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
            divider.line.fill.background()
            
            if slide_data['subtitle']:
                p2 = tf.add_paragraph()
                p2.text = slide_data['subtitle']
                p2.font.name = 'Microsoft JhengHei'
                p2.font.size = Pt(22)
                p2.font.color.rgb = RGBColor(2, 132, 199) # Sky-700
                p2.space_before = Pt(40)
                p2.alignment = PP_ALIGN.CENTER
                
        elif slide_data['has_dual_quotes']:
            # Render slide title at the top
            if slide_data['title']:
                # Accent bar
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
                accent_bar.line.fill.background()

                txBox = slide.shapes.add_textbox(Inches(1.05), Inches(0.6), Inches(11.48), Inches(0.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data['title']
                p.font.name = 'Microsoft JhengHei'
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950
                
                # Title Divider
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(226, 232, 240) # Slate-200
                line.line.fill.background()
                
            # Render Dual Columns (Light Cards)
            case_study_text = ""
            for item in slide_data.get('items', []):
                t = item['text']
                if '實戰案例對撞' in t or '實戰案例' in t:
                    case_study_text = t
                    break
                    
            card_height = Inches(3.8) if case_study_text else Inches(4.8)
            
            # Left Card
            left_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(1.85), Inches(5.6), card_height
            )
            left_card.fill.solid()
            left_card.fill.fore_color.rgb = RGBColor(240, 249, 255) # Sky-50
            left_card.line.color.rgb = RGBColor(186, 230, 253) # Sky-200
            left_card.line.width = Pt(1.5)
            
            tf_left = left_card.text_frame
            tf_left.word_wrap = True
            tf_left.margin_left = Inches(0.25)
            tf_left.margin_right = Inches(0.25)
            tf_left.margin_top = Inches(0.3)
            tf_left.margin_bottom = Inches(0.3)
            
            quotes = slide_data['quotes_data']
            
            p_left = tf_left.paragraphs[0]
            p_left.text = quotes[0]['name']
            p_left.font.name = 'Microsoft JhengHei'
            p_left.font.size = Pt(20)
            p_left.font.bold = True
            p_left.font.color.rgb = RGBColor(3, 105, 161) # Sky-700
            p_left.space_after = Pt(12)
            
            p_left_body = tf_left.add_paragraph()
            add_bold_text_run(p_left_body, quotes[0]['quote'], RGBColor(71, 85, 105), RGBColor(3, 105, 161))
            p_left_body.font.name = 'Microsoft JhengHei'
            p_left_body.font.size = Pt(15)
            
            # Right Card
            right_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.8), Inches(1.85), Inches(5.7), card_height
            )
            right_card.fill.solid()
            right_card.fill.fore_color.rgb = RGBColor(248, 250, 252) # Slate-50
            right_card.line.color.rgb = RGBColor(226, 232, 240) # Slate-200
            right_card.line.width = Pt(1.5)
            
            tf_right = right_card.text_frame
            tf_right.word_wrap = True
            tf_right.margin_left = Inches(0.25)
            tf_right.margin_right = Inches(0.25)
            tf_right.margin_top = Inches(0.3)
            tf_right.margin_bottom = Inches(0.3)
            
            p_right = tf_right.paragraphs[0]
            p_right.text = quotes[1]['name'] if len(quotes) > 1 else "東方聖人"
            p_right.font.name = 'Microsoft JhengHei'
            p_right.font.size = Pt(20)
            p_right.font.bold = True
            p_right.font.color.rgb = RGBColor(71, 85, 105) # Slate-600
            p_right.space_after = Pt(12)
            
            p_right_body = tf_right.add_paragraph()
            r_quote = quotes[1]['quote'] if len(quotes) > 1 else ""
            add_bold_text_run(p_right_body, r_quote, RGBColor(71, 85, 105), RGBColor(71, 85, 105))
            p_right_body.font.name = 'Microsoft JhengHei'
            p_right_body.font.size = Pt(15)
            
            # Render bottom case study callout if present
            if case_study_text:
                draw_callout_box(
                    slide, 
                    case_study_text, 
                    Inches(5.85), 
                    Inches(0.95),
                    RGBColor(255, 251, 235), # Amber-50 bg
                    RGBColor(253, 230, 138), # Amber-200 border
                    RGBColor(180, 83, 9)     # Amber-700 text accent
                )
            
        elif slide_data['table_data']:
            # Render Table Slide
            if slide_data['title']:
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
                accent_bar.line.fill.background()

                txBox = slide.shapes.add_textbox(Inches(1.05), Inches(0.6), Inches(11.48), Inches(0.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data['title']
                p.font.name = 'Microsoft JhengHei'
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950
                
                # Horizontal line under title
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(226, 232, 240)
                line.line.fill.background()
                
            headers = slide_data['table_data']['headers']
            rows = slide_data['table_data']['rows']
            
            rows_count = len(rows) + 1
            cols_count = len(headers)
            
            left = Inches(0.8)
            top = Inches(1.85)
            width = Inches(11.73)
            row_height = Inches(0.5)
            height = row_height * rows_count
            if height > Inches(4.8):
                height = Inches(4.8)
                
            table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
            table = table_shape.table
            
            col_width = int(width / cols_count)
            for col in table.columns:
                col.width = col_width
                
            for c_idx, h_text in enumerate(headers):
                cell = table.cell(0, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(8, 51, 68) # Cyan-950 Deep Water Blue
                
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = h_text
                p.font.name = 'Microsoft JhengHei'
                p.font.size = Pt(15)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
                
            for r_idx, row_cells in enumerate(rows):
                bg_color = RGBColor(248, 250, 252) if r_idx % 2 == 0 else RGBColor(255, 255, 255)
                for c_idx, cell_text in enumerate(row_cells):
                    cell = table.cell(r_idx + 1, c_idx)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg_color
                    
                    tf = cell.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    add_bold_text_run(p, cell_text, RGBColor(51, 65, 85), RGBColor(3, 105, 161))
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(13)
                    p.alignment = PP_ALIGN.CENTER if c_idx == 0 else PP_ALIGN.LEFT
                    
        else:
            # Standard Content Slide
            if slide_data['title']:
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
                accent_bar.line.fill.background()

                txBox = slide.shapes.add_textbox(Inches(1.05), Inches(0.6), Inches(11.48), Inches(0.8))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data['title']
                p.font.name = 'Microsoft JhengHei'
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950
                
                # Title Divider
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(226, 232, 240)
                line.line.fill.background()
                
            items = slide_data['items']
            
            # Find and separate Callout Boxes:
            # 1. 觀念核心
            core_concept_text = ""
            # 2. 實戰案例對撞
            case_study_text = ""
            
            core_items = []
            for item in items:
                t = item['text']
                if '觀念核心：' in t or '觀念核心' in t:
                    core_concept_text = t
                elif '實戰案例對撞：' in t or '實戰案例對撞' in t:
                    case_study_text = t
                else:
                    core_items.append(item)
                    
            top_offset = Inches(1.65)
            content_height = Inches(5.0)
            
            if core_concept_text:
                # Draw top callout box
                draw_callout_box(
                    slide, 
                    core_concept_text, 
                    top_offset, 
                    Inches(0.85),
                    RGBColor(240, 249, 255), # Sky-50 bg
                    RGBColor(186, 230, 253), # Sky-200 border
                    RGBColor(2, 132, 199)    # Sky-700 text accent
                )
                top_offset += Inches(1.0)
                content_height -= Inches(1.0)
                
            if case_study_text:
                # Draw bottom callout box
                draw_callout_box(
                    slide, 
                    case_study_text, 
                    Inches(6.0), 
                    Inches(0.9),
                    RGBColor(255, 251, 235), # Amber-50 bg
                    RGBColor(253, 230, 138), # Amber-200 border
                    RGBColor(180, 83, 9)     # Amber-700 text accent
                )
                content_height -= Inches(1.0)
                
            # Render remaining items in list format
            if core_items:
                contentBox = slide.shapes.add_textbox(Inches(0.8), top_offset, Inches(11.73), content_height)
                tf = contentBox.text_frame
                tf.word_wrap = True
                
                first_item = True
                for item in core_items:
                    if first_item:
                        p = tf.paragraphs[0]
                        first_item = False
                    else:
                        p = tf.add_paragraph()
                        
                    if item['level'] == 0:
                        p.space_before = Pt(8)
                        p.space_after = Pt(12)
                        
                        run_bullet = p.add_run()
                        run_bullet.text = "✦  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(28)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(14, 165, 233) # Sky-500
                    else:
                        p.space_before = Pt(4)
                        p.space_after = Pt(8)
                        p.left_indent = Inches(0.5)
                        
                        run_bullet = p.add_run()
                        run_bullet.text = "•  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(25)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(125, 211, 252) # Sky-300
                        
                    text_content = item['text']
                    if item['type'] == 'num_list':
                        text_content = f"{item['num']}. {text_content}"
                        
                    add_bold_text_run(p, text_content, RGBColor(51, 65, 85), RGBColor(3, 105, 161))
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(28 - item['level'] * 3)

    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python md_to_pptx_jiuyang.py <input_md_file> <output_pptx_file>")
        sys.exit(1)
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        
    slides_data = parse_jiuyang_markdown(input_file)
    if slides_data:
        # First split content slides
        slides_data = split_slides(slides_data)
        # Then split long table slides
        slides_data = split_table_slides(slides_data)
        create_presentation(slides_data, output_file)
