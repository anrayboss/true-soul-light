import re
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_markdown_to_slides(filepath):
    if not os.path.exists(filepath):
        print(f"Error: File not found: {filepath}")
        return []

    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    lines = content.split('\n')
    blocks = []
    current_block = []
    
    for line in lines:
        stripped = line.strip()
        
        is_separator = re.match(r'^\s*---\s*$', stripped)
        is_h1 = stripped.startswith('#') and not stripped.startswith('##')
        is_h2 = stripped.startswith('##') and not stripped.startswith('###')
        is_h3 = stripped.startswith('###') and not stripped.startswith('####')
        
        if is_separator or is_h1 or is_h2 or is_h3:
            if current_block:
                blocks.append("\n".join(current_block))
                current_block = []
            if not is_separator:
                current_block.append(line)
        else:
            current_block.append(line)
            
    if current_block:
        blocks.append("\n".join(current_block))
            
    slides = []
    
    for block in blocks:
        block = block.strip()
        if not block:
            continue
            
        # Parse notes (comments)
        notes = []
        html_notes = re.findall(r'<!--\s*Note:\s*([\s\S]*?)\s*-->', block, re.IGNORECASE)
        for n in html_notes:
            notes.append(n.strip())
        
        # Remove comments from slide body
        body = re.sub(r'<!--\s*Note:[\s\S]*?-->', '', block, flags=re.IGNORECASE).strip()
        
        # Match standalone line Notes:
        text_notes = re.findall(r'(?:^|\n)\s*Note:\s*([\s\S]*)$', body, re.IGNORECASE)
        for n in text_notes:
            notes.append(n.strip())
            
        body = re.sub(r'(?:^|\n)\s*Note:\s*[\s\S]*$', '', body, flags=re.IGNORECASE).strip()
        
        # Parse slide elements (title, lists, text)
        lines = body.split('\n')
        slide_title = ""
        slide_subtitle = ""
        items = []
        is_cover = False
        
        # Determine if it is a cover slide (starts with H1 '#')
        non_empty_lines = [l.strip() for l in lines if l.strip()]
        first_line = non_empty_lines[0] if non_empty_lines else ""
        
        if first_line.startswith('#') and not first_line.startswith('##'):
            slide_title = first_line[1:].strip()
            is_cover = True
            # Find subtitle if any
            for line in non_empty_lines[1:]:
                if line.startswith('##') and not line.startswith('###'):
                    slide_subtitle = line[2:].strip()
                    break
                elif line.startswith('- ') or line.startswith('* '):
                    slide_subtitle = line.strip()[2:]
                    break
        else:
            # Content slide parsing
            for line in lines:
                original_line = line
                line_stripped = line.strip()
                if not line_stripped:
                    continue
                
                # Check indentation before stripping
                leading_spaces = len(original_line) - len(original_line.lstrip())
                level = 1 if leading_spaces >= 2 else 0
                
                if line_stripped.startswith('##') and not line_stripped.startswith('###'):
                    slide_title = line_stripped[2:].strip()
                elif line_stripped.startswith('###') and not line_stripped.startswith('####'):
                    slide_title = line_stripped[3:].strip()
                elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                    items.append({'type': 'bullet', 'text': line_stripped[2:], 'level': level})
                elif re.match(r'^\d+\.\s+', line_stripped):
                    match = re.match(r'^(\d+)\.\s+(.*)', line_stripped)
                    items.append({'type': 'num_list', 'text': match.group(2), 'num': match.group(1), 'level': level})
                else:
                    items.append({'type': 'text', 'text': line_stripped, 'level': level})
                
        slides.append({
            'title': slide_title,
            'subtitle': slide_subtitle,
            'is_cover': is_cover,
            'items': items,
            'notes': "\n".join(notes)
        })
        
    return slides

def create_presentation(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply premium light background (#FCFCFC)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(252, 252, 252) 
        
        # Draw slide border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7)
        )
        border.fill.background() # Hollow fill
        border.line.color.rgb = RGBColor(226, 232, 240) # Slate-200
        border.line.width = Pt(1.0)
        
        # Add speaker notes if any
        if slide_data['notes']:
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
            
        if slide_data['is_cover']:
            # Create a large Title Box in the center
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 41, 59) # Slate-800
            p.alignment = PP_ALIGN.CENTER
            
            # Decorative centered line
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.03)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500 Accent Divider Line
            divider.line.fill.background()
            
            if slide_data['subtitle']:
                p2 = tf.add_paragraph()
                p2.text = slide_data['subtitle']
                p2.font.name = 'Microsoft JhengHei'
                p2.font.size = Pt(22)
                p2.font.color.rgb = RGBColor(2, 132, 199) # Sky-700 Accent Subtitle
                p2.space_before = Pt(40) # Push below divider
                p2.alignment = PP_ALIGN.CENTER
                
        else:
            # Content slide
            if not slide_data['items']:
                # Transition / Section Divider slide
                if slide_data['title']:
                    # Draw a nice clean vertical accent line in the middle
                    accent_bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(6.5), Inches(2.2), Inches(0.33), Inches(0.04)
                    )
                    accent_bar.fill.solid()
                    accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500 Accent Bar
                    accent_bar.line.fill.background()

                    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.333), Inches(2.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_data['title']
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(38)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950
                    p.alignment = PP_ALIGN.CENTER
            else:
                # Content slide with list items
                # 1. Slide Title (top left with modern vertical accent bar)
                if slide_data['title']:
                    # Modern vertical accent bar on the left of the title
                    accent_bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
                    )
                    accent_bar.fill.solid()
                    accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500 Accent Bar
                    accent_bar.line.fill.background()

                    txBox = slide.shapes.add_textbox(Inches(1.05), Inches(0.6), Inches(11.48), Inches(0.8))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = slide_data['title']
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950 / Deep Water Blue
                    
                    # Horizontal line under title
                    line = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
                    )
                    line.fill.solid()
                    line.fill.fore_color.rgb = RGBColor(226, 232, 240) # Slate-200 divider
                    line.line.fill.background()
                    
                # 2. Slide Content (middle)
                contentBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(11.73), Inches(4.8))
                tf = contentBox.text_frame
                tf.word_wrap = True
                
                first_item = True

                for item in slide_data['items']:
                    if first_item:
                        p = tf.paragraphs[0]
                        first_item = False
                    else:
                        p = tf.add_paragraph()
                        
                    # Custom spacing and bullet logic
                    if item['level'] == 0:
                        p.space_before = Pt(8)
                        p.space_after = Pt(12)
                        
                        # Star Bullet for level 0
                        run_bullet = p.add_run()
                        run_bullet.text = "✦  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(28)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(14, 165, 233) # Sky-500 Star Bullet
                    else:
                        p.space_before = Pt(4)
                        p.space_after = Pt(8)
                        p.left_indent = Inches(0.5)
                        
                        # Dot Bullet for level 1
                        run_bullet = p.add_run()
                        run_bullet.text = "•  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(25)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(125, 211, 252) # Sky-300 Dot Bullet
                    
                    text_content = item['text']
                    if item['type'] == 'num_list':
                        text_content = f"{item['num']}. {text_content}"
                        
                    # Support **bold** formatting runs inside text
                    parts = re.split(r'(\*\*.*?\*\*)', text_content)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run()
                            run.text = part[2:-2]
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(3, 105, 161) # Sky-700 bold text
                        else:
                            run = p.add_run()
                            run.text = part
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(51, 65, 85) # Slate-700 regular text
                            
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(28 - item['level'] * 3)
                    
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

def split_slides(slides_data):
    new_slides = []
    for slide in slides_data:
        if slide['is_cover'] or not slide['items']:
            new_slides.append(slide)
            continue
        
        items = slide['items']
        chunks = []
        current_chunk = []
        current_lines = 0
        
        for item in items:
            text = item['text']
            clean_text = text.replace('**', '')
            # Level 0 font size is 28pt, Level 1 is 25pt. Estimate characters per line:
            char_limit = 43 if item['level'] > 0 else 38
            item_lines = max(1, (len(clean_text) + char_limit - 1) // char_limit)
            
            if current_lines + item_lines > 10 and current_chunk:
                chunks.append(current_chunk)
                current_chunk = [item]
                current_lines = item_lines
            else:
                current_chunk.append(item)
                current_lines += item_lines
                
        if current_chunk:
            chunks.append(current_chunk)
            
        total_pages = len(chunks)
        if total_pages <= 1:
            new_slides.append(slide)
            continue
            
        for idx, chunk in enumerate(chunks):
            page_num = idx + 1
            original_title = slide['title']
            new_title = f"{original_title} ({page_num}/{total_pages})" if original_title else f"({page_num}/{total_pages})"
            
            new_slides.append({
                'title': new_title,
                'subtitle': slide['subtitle'],
                'is_cover': slide['is_cover'],
                'items': chunk,
                'notes': slide['notes']
            })
            
    return new_slides


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python md_to_pptx_white.py <input_md_file> <output_pptx_file>")
        sys.exit(1)
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        
    slides_data = parse_markdown_to_slides(input_file)
    if slides_data:
        slides_data = split_slides(slides_data)
        create_presentation(slides_data, output_file)
