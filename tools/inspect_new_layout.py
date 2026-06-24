import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

prs = Presentation("最終輸出用檔案/說明會新排版嘗試.pptx")
print(f"Total slides: {len(prs.slides)}")

for idx in range(len(prs.slides)):
    slide = prs.slides[idx]
    print(f"\n--- Slide {idx+1} ---")
    
    # Check if there is slide title
    title = ""
    if slide.shapes.title:
        title = slide.shapes.title.text
        print(f"Slide Title: '{title}'")
        
    for shape in slide.shapes:
        print(f"Shape: '{shape.name}', Type: {shape.shape_type}")
        left = shape.left.inches if shape.left else 0
        top = shape.top.inches if shape.top else 0
        width = shape.width.inches if shape.width else 0
        height = shape.height.inches if shape.height else 0
        print(f"  Pos: Left={left:.3f}\", Top={top:.3f}\", Width={width:.3f}\", Height={height:.3f}\"")
        
        # Check text
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, p in enumerate(tf.paragraphs):
                if p.text.strip():
                    print(f"    Paragraph {p_idx+1}: '{p.text[:120]}' (Alignment: {p.alignment})")
                    for r_idx, r in enumerate(p.runs):
                        size_str = f"{r.font.size.pt}pt" if r.font.size else "None"
                        print(f"      Run {r_idx+1}: '{r.text[:40]}' | Font: {r.font.name}, Size: {size_str}, Bold: {r.font.bold}, Italic: {r.font.italic}")
