import os
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def extract_19_steps():
    teaching_path = "最終輸出用檔案/陳修平短影音保姆級教學.pptx"
    fallback_steps = [
        "✦ 1. 調亮螢幕：拍攝前如發現觀景螢幕偏暗，請先手動將螢幕亮度調亮，以利準確檢視畫面細節。",
        "✦ 2. 改設 60 幀：將左上角的參數設定（如 HD 30 幀）全部改為 60 幀（60 fps），拍出來的畫面更順暢、沒有延遲感。",
        "✦ 3. 自拍用前鏡頭：自己拍自己時可直接用前鏡頭，雖然畫質稍遜，但只要陽光等光線充足，畫面效果就不受影響。",
        "✦ 4. 拍人固定 1x 焦距：拍攝他人時請固定使用 1 倍（1x）焦距，切勿使用廣角，以免人像邊緣變形。",
        "✦ 5. 拍攝者自行移動：需要移動畫面時，請拍攝者自己移動，不要叫主角（被攝者）移動。",
        "✦ 6. 景物直接靠過去：拍風景或樹木時，直接往前走靠過去拍，不要用雙指在螢幕上拉近拉遠變焦。",
        "✦ 7. 廣角（0.5x）保持水平：使用 0.5 倍超廣角拍攝時必須保持手機水平。稍微傾斜會讓主角顯得極矮或極胖。",
        "✦ 8. 水平同步跟拍：主角在行進或移動時，拍攝者請維持與主角水平同步的速度跟著走，畫面才會穩定。",
        "✦ 9. 開啟相機隔線：進入手機設定將相機的「格線/隔線」輔助功能打開，便於視覺定位。",
        "✦ 10. 主角置中與額頭對齊：構圖時將主角放在畫面正中間，並讓額頭對齊格線最上方的那條線，比例最為美觀。",
        "✦ 11. 避開穿頭殺（穿駝紗）：注意被攝者講話時，頭頂後方不要有柱子、樹木或其他雜物穿過，避免干擾聽眾的視覺專注度。",
        "✦ 12. 側邊 45 度角拍攝：幫別人拍攝時，不要直接擋在對方正前方（會給被攝者極大壓力且阻擋前進），應站在 45 度角斜前方拍攝。",
        "✦ 13. 對視引導法（鏡頭避讓）：將手機稍微移到側邊，讓被攝者看著拍攝者的眼睛說話。雖然他是在跟拍攝者對話，但在畫面上看起來就像是自然地對著鏡頭說話，能有效降低被攝者的鏡頭壓力。",
        "✦ 14. 句末停頓 1 秒：念稿或說話時，每講完一句話先停頓 1 秒，以利後製剪輯時切片與淡入淡出。",
        "✦ 15. 小指當底座支撐：手持手機時，用小手指墊在手機底部當作支架。有了支撐點，手持拍攝時間變長且不易手震。",
        "✦ 16. 手指避開鏡頭：現代手機鏡頭越來越多，握持時務必留意手指不要遮擋住任何一顆鏡頭。",
        "✦ 17. 背景避開大白牆：拍攝口播影片時，背景最好不要是單調的大白牆（視覺單調且無法說明所處情境），建議在戶外或畫面有延伸感的環境中拍攝。",
        "✦ 18. 外接麥克風必先測音：如使用外接式指向麥克風，錄製前請務必先錄製一小段並播放聆聽，確認收音正常再正式錄製。",
        "✦ 19. 一鏡到底不中斷：不論講錯、吃螺絲或頭髮亂了，都不要關掉相機重新開錄。維持一鏡到底錄成一個長檔案。雖然檔案可能很大（如 20 分鐘），但剪輯時只有一個軌道，比整理數十個零碎檔案更容易進行切片與後製。"
    ]
    
    if not os.path.exists(teaching_path):
        print(f"Teaching PPTX not found, using built-in steps.")
        return fallback_steps
        
    try:
        prs = Presentation(teaching_path)
        extracted = []
        # Slides index 2 to 7 (Slides 3 to 8)
        for idx in range(2, len(prs.slides)):
            slide = prs.slides[idx]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        text = p.text.strip()
                        if not text:
                            continue
                        # Split by slash or newlines
                        parts = re.split(r'\n|(?<=.)\s+/\s+|(?<=.)\s+／\s+', text)
                        for part in parts:
                            part = part.strip()
                            if not part:
                                continue
                            if part.startswith("✦") or re.match(r'^(?:✦\s*)?\d+', part):
                                extracted.append(part)
        if len(extracted) >= 15:
            print(f"Successfully extracted {len(extracted)} steps from teaching slides.")
            return extracted
        else:
            print(f"Extracted steps too few ({len(extracted)}), falling back to robust built-in steps.")
            return fallback_steps
    except Exception as e:
        print(f"Error reading teaching slides: {e}. Using fallback steps.")
        return fallback_steps

def process_shape_recursive(shape, title_shape, is_main_list, width):
    if shape == title_shape:
        return
        
    # Recurse if group shape
    if shape.shape_type == 6 or hasattr(shape, "shapes"): # MSO_SHAPE_TYPE.GROUP is 6
        try:
            for subshape in shape.shapes:
                sub_width = subshape.width.inches if subshape.width else 0
                process_shape_recursive(subshape, title_shape, is_main_list, sub_width)
        except Exception as e:
            print(f"Error recursing group shape: {e}")
        return
        
    # Table processing
    if shape.has_table:
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.has_text_frame:
                        for p in cell.text_frame.paragraphs:
                            p.font.name = "微軟正黑體"
                            for run in p.runs:
                                run.font.name = "微軟正黑體"
        except Exception as e:
            print(f"Error processing table: {e}")
            
    # Process text frames
    if shape.has_text_frame:
        tf = shape.text_frame
        tf.word_wrap = True
        for p in tf.paragraphs:
            p.font.name = "微軟正黑體"
            for run in p.runs:
                run.font.name = "微軟正黑體"
                
                # Determine font size
                if is_main_list:
                    run.font.size = Pt(20) # 20pt for main list text
                else:
                    orig_size = run.font.size
                    if orig_size:
                        run.font.size = orig_size
                    else:
                        text_len = len(p.text)
                        if text_len < 6 and width < 2.0:
                            run.font.size = Pt(18)
                        elif width < 1.5:
                            run.font.size = Pt(11)
                        elif width < 3.0:
                            run.font.size = Pt(13)
                        else:
                            run.font.size = Pt(16)
                            
    # Apply style tuning to shapes
    if shape.shape_type in [MSO_SHAPE.OVAL, MSO_SHAPE.ROUNDED_RECTANGLE, MSO_SHAPE.RECTANGLE]:
        # If shape has fill, set to premium soft sky blue
        try:
            if shape.fill.type == 1: # Solid fill
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 249, 255) # Premium Sky-50
            if hasattr(shape, 'line') and shape.line:
                shape.line.color.rgb = RGBColor(186, 230, 253) # Premium Sky-200
                shape.line.width = Pt(1.5)
        except Exception:
            pass
    elif shape.shape_type in [MSO_SHAPE.DOWN_ARROW, MSO_SHAPE.UP_ARROW, MSO_SHAPE.LEFT_ARROW, MSO_SHAPE.RIGHT_ARROW, MSO_SHAPE.UP_DOWN_ARROW]:
        # Style arrows with Sky-500
        try:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
            if hasattr(shape, 'line') and shape.line:
                shape.line.color.rgb = RGBColor(14, 165, 233)
        except Exception:
            pass

def main():
    input_path = "最終輸出用檔案/圖表重製.pptx"
    output_path = "最終輸出用檔案/圖表重製_排版美編完成.pptx"
    
    if not os.path.exists(input_path):
        print(f"Error: Input presentation not found: {input_path}")
        sys.exit(1)
        
    # Read the 19 steps
    steps_19 = extract_19_steps()
    
    prs = Presentation(input_path)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print(f"Processing {len(prs.slides)} slides from template...")
    
    for slide_idx, slide in enumerate(prs.slides):
        # 1. Apply premium background color (#FCFCFC)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(252, 252, 252)
        
        # 2. Identify and remove outer slide border shape if exists
        shapes_to_delete = []
        for shape in slide.shapes:
            # Border shape: Rectangle, left & top around 0.4", width > 12", height > 6.5", no text
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                w = shape.width.inches if shape.width else 0
                h = shape.height.inches if shape.height else 0
                if left < 0.6 and top < 0.6 and w > 12.0 and h > 6.5:
                    if not shape.has_text_frame or not any(p.text.strip() for p in shape.text_frame.paragraphs):
                        shapes_to_delete.append(shape)
                        
        for shape in shapes_to_delete:
            try:
                shape.element.getparent().remove(shape.element)
                print(f"Removed outer border on Slide {slide_idx+1}")
            except Exception as e:
                print(f"Failed to remove border: {e}")
                
        # 3. Detect and style slide title
        title_shape = None
        # Try finding title placeholder first
        if slide.shapes.title:
            title_shape = slide.shapes.title
        else:
            # Fallback: look for a top textbox that looks like a title
            for shape in slide.shapes:
                if shape.has_text_frame:
                    top = shape.top.inches if shape.top else 99
                    h = shape.height.inches if shape.height else 0
                    text = "".join(p.text for p in shape.text_frame.paragraphs).strip()
                    if top < 1.4 and len(text) > 0 and len(text) < 50:
                        title_shape = shape
                        break
                        
        has_title = False
        if title_shape and title_shape.has_text_frame:
            title_text = "".join(p.text for p in title_shape.text_frame.paragraphs).strip()
            if title_text:
                has_title = True
                # Move to new standard title position
                title_shape.left = Inches(1.05)
                title_shape.top = Inches(0.60)
                title_shape.width = Inches(11.48)
                title_shape.height = Inches(0.8)
                
                tf = title_shape.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                
                # Format paragraphs
                for p in tf.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    for run in p.runs:
                        run.font.name = "微軟正黑體"
                        run.font.size = Pt(32)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(8, 51, 68) # Cyan-950
                
                # Draw title vertical Accent Bar
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-500
                accent_bar.line.fill.background()
                
                # Draw Title Divider Line
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(226, 232, 240) # Slate-200
                line.line.fill.background()
                
        # 4a. Handle Slide 10 - Reconstruct with 10 Trans-codes and 5 Exploding templates
        # Note: Slide index is 9 (0-indexed)
        if slide_idx == 9:
            # Delete all shapes except title, Accent Bar and Divider Line
            to_delete = []
            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                if left == 0.8 and (top == 0.65 or top == 1.45):
                    continue
                to_delete.append(shape)
                
            for shape in to_delete:
                try:
                    shape.element.getparent().remove(shape.element)
                except Exception:
                    pass
            
            # Recreate title
            title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.60), Inches(11.48), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "推廣 6 —— 客從哪來（What）：10 大轉單密碼 ╳ 5 大爆文套路"
            p.font.name = "微軟正黑體"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(8, 51, 68)
            
            # Add Accent bar and Divider
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(14, 165, 233)
            accent.line.fill.background()
            
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01))
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(226, 232, 240)
            div.line.fill.background()
            
            # Left Column Header
            left_header = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(0.8))
            tf_lh = left_header.text_frame
            tf_lh.word_wrap = True
            p_lh = tf_lh.paragraphs[0]
            p_lh.text = "10 大轉單密碼"
            p_lh.font.name = "微軟正黑體"
            p_lh.font.size = Pt(24)
            p_lh.font.bold = True
            p_lh.font.color.rgb = RGBColor(3, 105, 161) # Sky-700
            
            p_lh_sub = tf_lh.add_paragraph()
            p_lh_sub.text = "撬開消費者的心理防線"
            p_lh_sub.font.name = "微軟正黑體"
            p_lh_sub.font.size = Pt(14)
            p_lh_sub.font.color.rgb = RGBColor(100, 116, 139) # Slate-500
            
            # Right Column Header
            right_header = slide.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(4.8), Inches(0.8))
            tf_rh = right_header.text_frame
            tf_rh.word_wrap = True
            p_rh = tf_rh.paragraphs[0]
            p_rh.text = "5 大爆文套路"
            p_rh.font.name = "微軟正黑體"
            p_rh.font.size = Pt(24)
            p_rh.font.bold = True
            p_rh.font.color.rgb = RGBColor(3, 105, 161) # Sky-700
            
            p_rh_sub = tf_rh.add_paragraph()
            p_rh_sub.text = "駕馭流量的爆文套路"
            p_rh_sub.font.name = "微軟正黑體"
            p_rh_sub.font.size = Pt(14)
            p_rh_sub.font.color.rgb = RGBColor(100, 116, 139) # Slate-500
            
            # Add Left Cards (10 items, 2 cols x 5 rows)
            left_items = [
                "1. 獨家賣點", "6. 成功見證",
                "2. 急迫性 / 稀缺性", "7. 踩痛點",
                "3. 談錢", "8. 貼標籤",
                "4. 掛保證", "9. 給利益",
                "5. 立權威", "10. 隱藏的第二標題 (P.S.)"
            ]
            
            card_tops = [2.6, 3.4, 4.2, 5.0, 5.8]
            
            for row_idx in range(5):
                top_pos = card_tops[row_idx]
                
                # Left Column Card (Col 1)
                item1 = left_items[row_idx * 2]
                card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(2.8), Inches(0.65))
                card1.fill.solid()
                card1.fill.fore_color.rgb = RGBColor(240, 249, 255) # Sky-50
                card1.line.color.rgb = RGBColor(186, 230, 253) # Sky-200
                card1.line.width = Pt(1.5)
                tf1 = card1.text_frame
                tf1.word_wrap = True
                p1 = tf1.paragraphs[0]
                p1.text = item1
                p1.alignment = PP_ALIGN.CENTER
                p1.font.name = "微軟正黑體"
                p1.font.size = Pt(15)
                p1.font.bold = True
                p1.font.color.rgb = RGBColor(3, 105, 161)
                
                # Right Column Card (Col 2)
                item2 = left_items[row_idx * 2 + 1]
                card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(top_pos), Inches(2.8), Inches(0.65))
                card2.fill.solid()
                card2.fill.fore_color.rgb = RGBColor(240, 249, 255)
                card2.line.color.rgb = RGBColor(186, 230, 253)
                card2.line.width = Pt(1.5)
                tf2 = card2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = item2
                p2.alignment = PP_ALIGN.CENTER
                p2.font.name = "微軟正黑體"
                p2.font.size = Pt(15)
                p2.font.bold = True
                p2.font.color.rgb = RGBColor(3, 105, 161)
                
            # Add Right Cards (5 items, 1 col x 5 rows)
            right_items = [
                "1. 列點文",
                "2. 懶人包",
                "3. 教學文",
                "4. 搬知識",
                "5. 說故事"
            ]
            for row_idx in range(5):
                top_pos = card_tops[row_idx]
                item = right_items[row_idx]
                
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(top_pos), Inches(4.5), Inches(0.65))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(240, 249, 255)
                card.line.color.rgb = RGBColor(186, 230, 253)
                card.line.width = Pt(1.5)
                tf_c = card.text_frame
                tf_c.word_wrap = True
                p_c = tf_c.paragraphs[0]
                p_c.text = item
                p_c.alignment = PP_ALIGN.CENTER
                p_c.font.name = "微軟正黑體"
                p_c.font.size = Pt(15)
                p_c.font.bold = True
                p_c.font.color.rgb = RGBColor(3, 105, 161)
                
            print("Successfully populated Slide 10 with 10 Trans-codes and 5 Exploding templates.")
            continue

        # 4b. Handle Slide 11 - Reconstruct with 7 Content Themes and 11 Traffic Codes
        # Note: Slide index is 10 (0-indexed)
        if slide_idx == 10:
            # Delete all shapes except title, Accent Bar and Divider Line
            to_delete = []
            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                if left == 0.8 and (top == 0.65 or top == 1.45):
                    continue
                to_delete.append(shape)
                
            for shape in to_delete:
                try:
                    shape.element.getparent().remove(shape.element)
                except Exception:
                    pass
            
            # Recreate title
            title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.60), Inches(11.48), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = "推廣 6 —— 客從哪來（What）：7 大內容主題 ╳ 11 大流量密碼"
            p.font.name = "微軟正黑體"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(8, 51, 68)
            
            # Add Accent bar and Divider
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(14, 165, 233)
            accent.line.fill.background()
            
            div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01))
            div.fill.solid()
            div.fill.fore_color.rgb = RGBColor(226, 232, 240)
            div.line.fill.background()
            
            # Left Column Header
            left_header = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(7.2), Inches(0.8))
            tf_lh = left_header.text_frame
            tf_lh.word_wrap = True
            p_lh = tf_lh.paragraphs[0]
            p_lh.text = "7 大內容主題"
            p_lh.font.name = "微軟正黑體"
            p_lh.font.size = Pt(24)
            p_lh.font.bold = True
            p_lh.font.color.rgb = RGBColor(3, 105, 161)
            
            p_lh_sub = tf_lh.add_paragraph()
            p_lh_sub.text = "打造社群多元價值的內容框架"
            p_lh_sub.font.name = "微軟正黑體"
            p_lh_sub.font.size = Pt(14)
            p_lh_sub.font.color.rgb = RGBColor(100, 116, 139)
            
            # Right Column Header
            right_header = slide.shapes.add_textbox(Inches(8.8), Inches(1.7), Inches(3.8), Inches(0.8))
            tf_rh = right_header.text_frame
            tf_rh.word_wrap = True
            p_rh = tf_rh.paragraphs[0]
            p_rh.text = "11 大流量密碼"
            p_rh.font.name = "微軟正黑體"
            p_rh.font.size = Pt(24)
            p_rh.font.bold = True
            p_rh.font.color.rgb = RGBColor(3, 105, 161)
            
            p_rh_sub = tf_rh.add_paragraph()
            p_rh_sub.text = "爆發流量的社群密碼"
            p_rh_sub.font.name = "微軟正黑體"
            p_rh_sub.font.size = Pt(14)
            p_rh_sub.font.color.rgb = RGBColor(100, 116, 139)
            
            # Left Section Cards
            # Sub-col 1 (Left 0.8", Width 3.6")
            # 1. 曬過程
            card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.6), Inches(3.6), Inches(1.2))
            card1.fill.solid()
            card1.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card1.line.color.rgb = RGBColor(186, 230, 253)
            card1.line.width = Pt(1.5)
            tf1 = card1.text_frame
            tf1.word_wrap = True
            p1_t = tf1.paragraphs[0]
            p1_t.text = "❶ 曬過程"
            p1_t.font.name = "微軟正黑體"
            p1_t.font.size = Pt(15)
            p1_t.font.bold = True
            p1_t.font.color.rgb = RGBColor(3, 105, 161)
            p1_b = tf1.add_paragraph()
            p1_b.text = "製作過程 · 改造計畫 · 產品測試\n體驗分享 · 任務挑戰 · 成就故事"
            p1_b.font.name = "微軟正黑體"
            p1_b.font.size = Pt(11)
            p1_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # 2. 說故事
            card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.95), Inches(3.6), Inches(0.8))
            card2.fill.solid()
            card2.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card2.line.color.rgb = RGBColor(186, 230, 253)
            card2.line.width = Pt(1.5)
            tf2 = card2.text_frame
            tf2.word_wrap = True
            p2_t = tf2.paragraphs[0]
            p2_t.text = "❷ 說故事"
            p2_t.font.name = "微軟正黑體"
            p2_t.font.size = Pt(15)
            p2_t.font.bold = True
            p2_t.font.color.rgb = RGBColor(3, 105, 161)
            p2_b = tf2.add_paragraph()
            p2_b.text = "案例故事 · 他人故事"
            p2_b.font.name = "微軟正黑體"
            p2_b.font.size = Pt(11)
            p2_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # 3. 教知識
            card3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(3.6), Inches(1.2))
            card3.fill.solid()
            card3.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card3.line.color.rgb = RGBColor(186, 230, 253)
            card3.line.width = Pt(1.5)
            tf3 = card3.text_frame
            tf3.word_wrap = True
            p3_t = tf3.paragraphs[0]
            p3_t.text = "❸ 教知識"
            p3_t.font.name = "微軟正黑體"
            p3_t.font.size = Pt(15)
            p3_t.font.bold = True
            p3_t.font.color.rgb = RGBColor(3, 105, 161)
            p3_b = tf3.add_paragraph()
            p3_b.text = "行業揭秘 · 乾貨輸出 · 資訊整理\n正向推薦 · 反向推薦"
            p3_b.font.name = "微軟正黑體"
            p3_b.font.size = Pt(11)
            p3_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # Sub-col 2 (Left 4.6", Width 3.6")
            # 4. 熬雞湯
            card4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(2.6), Inches(3.6), Inches(0.95))
            card4.fill.solid()
            card4.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card4.line.color.rgb = RGBColor(186, 230, 253)
            card4.line.width = Pt(1.5)
            tf4 = card4.text_frame
            tf4.word_wrap = True
            p4_t = tf4.paragraphs[0]
            p4_t.text = "❹ 熬雞湯"
            p4_t.font.name = "微軟正黑體"
            p4_t.font.size = Pt(15)
            p4_t.font.bold = True
            p4_t.font.color.rgb = RGBColor(3, 105, 161)
            p4_b = tf4.add_paragraph()
            p4_b.text = "自我成長 · 人際關係 · 處事之道"
            p4_b.font.name = "微軟正黑體"
            p4_b.font.size = Pt(11)
            p4_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # 5. 選立場
            card5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(3.7), Inches(3.6), Inches(0.95))
            card5.fill.solid()
            card5.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card5.line.color.rgb = RGBColor(186, 230, 253)
            card5.line.width = Pt(1.5)
            tf5 = card5.text_frame
            tf5.word_wrap = True
            p5_t = tf5.paragraphs[0]
            p5_t.text = "❺ 選立場"
            p5_t.font.name = "微軟正黑體"
            p5_t.font.size = Pt(15)
            p5_t.font.bold = True
            p5_t.font.color.rgb = RGBColor(3, 105, 161)
            p5_b = tf5.add_paragraph()
            p5_b.text = "正面觀點 · 反面觀點 · 中立觀點"
            p5_b.font.name = "微軟正黑體"
            p5_b.font.size = Pt(11)
            p5_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # 6. 演劇情
            card6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(4.8), Inches(3.6), Inches(0.95))
            card6.fill.solid()
            card6.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card6.line.color.rgb = RGBColor(186, 230, 253)
            card6.line.width = Pt(1.5)
            tf6 = card6.text_frame
            tf6.word_wrap = True
            p6_t = tf6.paragraphs[0]
            p6_t.text = "❻ 演劇情"
            p6_t.font.name = "微軟正黑體"
            p6_t.font.size = Pt(15)
            p6_t.font.bold = True
            p6_t.font.color.rgb = RGBColor(3, 105, 161)
            p6_b = tf6.add_paragraph()
            p6_b.text = "POV · 熱門梗 · 生活觀察"
            p6_b.font.name = "微軟正黑體"
            p6_b.font.size = Pt(11)
            p6_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # 7. 賣產品
            card7 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(5.9), Inches(3.6), Inches(0.8))
            card7.fill.solid()
            card7.fill.fore_color.rgb = RGBColor(240, 249, 255)
            card7.line.color.rgb = RGBColor(186, 230, 253)
            card7.line.width = Pt(1.5)
            tf7 = card7.text_frame
            tf7.word_wrap = True
            p7_t = tf7.paragraphs[0]
            p7_t.text = "❼ 賣產品"
            p7_t.font.name = "微軟正黑體"
            p7_t.font.size = Pt(15)
            p7_t.font.bold = True
            p7_t.font.color.rgb = RGBColor(3, 105, 161)
            p7_b = tf7.add_paragraph()
            p7_b.text = "展現服務 · 展現產品"
            p7_b.font.name = "微軟正黑體"
            p7_b.font.size = Pt(11)
            p7_b.font.color.rgb = RGBColor(71, 85, 105)
            
            # Right Column Items (11 items)
            flow_box = slide.shapes.add_textbox(Inches(8.8), Inches(2.5), Inches(3.8), Inches(4.5))
            tf_flow = flow_box.text_frame
            tf_flow.word_wrap = True
            tf_flow.margin_left = tf_flow.margin_right = tf_flow.margin_top = tf_flow.margin_bottom = 0
            
            flow_items = [
                "1. 成本 (Cost)", "2. 受眾 (Audience)", "3. 隨機 (Random)",
                "4. 反差 (Contrast)", "5. 貼金 (Status/Flexing)", "6. 互動 (Interaction)",
                "7. 熱點 (Hotspots)", "8. 懷舊 (Nostalgia)", "9. 吐槽 (Roasting)",
                "10. 賀爾蒙 (Hormones)", "11. 光練不說 (Action only)"
            ]
            
            for idx, item in enumerate(flow_items):
                p = tf_flow.paragraphs[0] if idx == 0 else tf_flow.add_paragraph()
                p.text = f"✦  {item}"
                p.space_after = Pt(6)
                p.font.name = "微軟正黑體"
                p.font.size = Pt(13.5)
                p.font.bold = True
                p.font.color.rgb = RGBColor(71, 85, 105)
                
            print("Successfully populated Slide 11 with 7 Content Themes and 11 Traffic Codes.")
            continue

        # 4. Handle Slide 12 - Reconstruct with 19 Steps (Special Case)
        # Note: Slide index is 11 (0-indexed)
        if slide_idx == 11:
            # Delete all shapes except title, Accent Bar and Divider Line
            to_delete = []
            for shape in slide.shapes:
                if shape == title_shape:
                    continue
                # Keep newly added accent bar and divider line
                left = shape.left.inches if shape.left else 0
                top = shape.top.inches if shape.top else 0
                w = shape.width.inches if shape.width else 0
                h = shape.height.inches if shape.height else 0
                if left == 0.8 and (top == 0.65 or top == 1.45):
                    continue
                to_delete.append(shape)
                
            for shape in to_delete:
                try:
                    shape.element.getparent().remove(shape.element)
                except Exception as e:
                    pass
            
            # Recreate title if missing (just to be safe)
            if not has_title:
                # Add title
                title_box = slide.shapes.add_textbox(Inches(1.05), Inches(0.60), Inches(11.48), Inches(0.8))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = "推廣 6 —— 客從哪來（What）：短影音 19 步"
                p.font.name = "微軟正黑體"
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(8, 51, 68)
                
                # Add Accent bar and Divider
                accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.65), Inches(0.08), Inches(0.6))
                accent.fill.solid()
                accent.fill.fore_color.rgb = RGBColor(14, 165, 233)
                accent.line.fill.background()
                
                div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.73), Inches(0.01))
                div.fill.solid()
                div.fill.fore_color.rgb = RGBColor(226, 232, 240)
                div.line.fill.background()
                
            # Create Left Column (Steps 1-10)
            left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.2))
            tf_left = left_box.text_frame
            tf_left.word_wrap = True
            tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
            
            # Create Right Column (Steps 11-19)
            right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.75), Inches(5.6), Inches(5.2))
            tf_right = right_box.text_frame
            tf_right.word_wrap = True
            tf_right.margin_left = tf_right.margin_right = tf_right.margin_top = tf_right.margin_bottom = 0
            
            for idx, step_text in enumerate(steps_19):
                if idx < 10:
                    p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
                else:
                    p = tf_right.paragraphs[0] if idx == 10 else tf_right.add_paragraph()
                    
                p.space_after = Pt(6)
                
                # Color code step prefix
                parts = re.split(r'(✦\s*\d+\.\s*[^：]+：)', step_text)
                if len(parts) >= 3:
                    prefix = parts[1]
                    body = "".join(parts[2:])
                    
                    run_pre = p.add_run()
                    run_pre.text = prefix
                    run_pre.font.name = "微軟正黑體"
                    run_pre.font.size = Pt(11.5)
                    run_pre.font.bold = True
                    run_pre.font.color.rgb = RGBColor(3, 105, 161) # Sky-700
                    
                    run_body = p.add_run()
                    run_body.text = body
                    run_body.font.name = "微軟正黑體"
                    run_body.font.size = Pt(11.5)
                    run_body.font.color.rgb = RGBColor(71, 85, 105) # Slate-600
                else:
                    run = p.add_run()
                    run.text = step_text
                    run.font.name = "微軟正黑體"
                    run.font.size = Pt(11.5)
                    run.font.color.rgb = RGBColor(71, 85, 105)
            
            print("Successfully populated Slide 12 with 19 steps in double columns.")
            continue
            
        # 5. Format other shapes on non-special slides
        for shape in slide.shapes:
            if shape == title_shape:
                continue
                
            is_main_list = False
            left = shape.left.inches if shape.left else 0
            width = shape.width.inches if shape.width else 0
            if left < 2.5 and width > 6.0:
                is_main_list = True
                
            process_shape_recursive(shape, title_shape, is_main_list, width)
                    
    # Save the output file
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    try:
        prs.save(output_path)
        print(f"Reconstructed charts presentation saved to: {output_path}")
    except PermissionError:
        alt_path = output_path.replace(".pptx", "_v2.pptx")
        print(f"Permission denied on {output_path} (possibly open in PowerPoint). Saving to alternative path: {alt_path}")
        prs.save(alt_path)
        print(f"Reconstructed charts presentation saved to: {alt_path}")

if __name__ == "__main__":
    main()
