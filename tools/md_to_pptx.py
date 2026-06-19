import re
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_markdown_to_slides(filepath):
    if not os.path.exists(filepath):
        print(f"Error: File not found: {filepath}")
        return []

    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by horizontal rules
    blocks = re.split(r'\n\s*---\s*\n', content)
    slides = []
    
    for block in blocks:
        block = block.strip()
        if not block:
            continue
            
        # Parse notes (comments)
        notes = []
        html_notes = re.findall(r'<!--\s*Note:\s*([\s\S]*?)\s*-->', block, re.IGNORECASE)
        for n in html_notes:
            notes.append(n.strip())
        
        # Remove comments from slide body
        body = re.sub(r'<!--\s*Note:[\s\S]*?-->', '', block, flags=re.IGNORECASE).strip()
        
        # Match standalone line Notes:
        text_notes = re.findall(r'(?:^|\n)\s*Note:\s*([\s\S]*)$', body, re.IGNORECASE)
        for n in text_notes:
            notes.append(n.strip())
            
        body = re.sub(r'(?:^|\n)\s*Note:\s*[\s\S]*$', '', body, flags=re.IGNORECASE).strip()
        
        # Parse slide elements (title, lists, text)
        lines = body.split('\n')
        slide_title = ""
        slide_subtitle = ""
        items = []
        is_cover = False
        
        # Determine if it is a cover slide (first line is H1 '#')
        non_empty_lines = [l.strip() for l in lines if l.strip()]
        if non_empty_lines and non_empty_lines[0].startswith('# '):
            slide_title = non_empty_lines[0][2:]
            is_cover = True
            if len(non_empty_lines) > 1 and non_empty_lines[1].startswith('## '):
                slide_subtitle = non_empty_lines[1][3:]
        else:
            # Content slide parsing
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                if line.startswith('## ') or line.startswith('### '):
                    slide_title = line.split(maxsplit=1)[1]
                elif line.startswith('- ') or line.startswith('* '):
                    items.append({'type': 'bullet', 'text': line[2:], 'level': 0})
                elif line.startswith('  - ') or line.startswith('    - ') or line.startswith('  * ') or line.startswith('    * '):
                    items.append({'type': 'bullet', 'text': line.strip()[2:], 'level': 1})
                elif re.match(r'^\d+\.\s+', line):
                    match = re.match(r'^(\d+)\.\s+(.*)', line)
                    items.append({'type': 'num_list', 'text': match.group(2), 'num': match.group(1), 'level': 0})
                else:
                    items.append({'type': 'text', 'text': line, 'level': 0})
                
        slides.append({
            'title': slide_title,
            'subtitle': slide_subtitle,
            'is_cover': is_cover,
            'items': items,
            'notes': "\n".join(notes)
        })
        
    return slides

def create_presentation(slides_data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply dark background (Cyber Neon style: #090D16 deep slate background)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(9, 13, 22) 
        
        # Draw slide border (mimics the .slide-border in CSS)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.4), Inches(0.4), Inches(12.533), Inches(6.7)
        )
        border.fill.background() # Hollow fill
        border.line.color.rgb = RGBColor(3, 105, 161) # Soft dark sky-blue border line
        border.line.width = Pt(1.5)
        
        # Add speaker notes if any
        if slide_data['notes']:
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
            
        if slide_data['is_cover']:
            # Create a large Title Box in the center
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = slide_data['title']
            p.font.name = 'Microsoft JhengHei'
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255) # White title
            p.alignment = PP_ALIGN.CENTER
            
            # Decorative centered line
            divider = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(4.5), Inches(3.6), Inches(4.333), Inches(0.03)
            )
            divider.fill.solid()
            divider.fill.fore_color.rgb = RGBColor(14, 165, 233) # Sky-blue Accent Divider Line
            divider.line.fill.background()
            
            if slide_data['subtitle']:
                p2 = tf.add_paragraph()
                p2.text = slide_data['subtitle']
                p2.font.name = 'Microsoft JhengHei'
                p2.font.size = Pt(22)
                p2.font.color.rgb = RGBColor(14, 165, 233) # Sky-blue Accent Subtitle
                p2.space_before = Pt(40) # Push below divider
                p2.alignment = PP_ALIGN.CENTER
                
        else:
            # Content slide
            # 1. Slide Title (top left)
            if slide_data['title']:
                txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.2))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_data['title']
                p.font.name = 'Microsoft JhengHei'
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = RGBColor(14, 165, 233) # Sky-blue Title (#0EA5E9)
                
                # Horizontal line under title
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.02)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(49, 46, 129) # Subtle divider line
                line.line.fill.background()
                
            # 2. Slide Content (middle)
            if slide_data['items']:
                contentBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0))
                tf = contentBox.text_frame
                tf.word_wrap = True
                
                first_item = True
                for item in slide_data['items']:
                    if first_item:
                        p = tf.paragraphs[0]
                        first_item = False
                    else:
                        p = tf.add_paragraph()
                        
                    # Custom spacing and bullet logic instead of PPTX default bullets to replicate style
                    if item['level'] == 0:
                        p.space_after = Pt(12)
                        
                        # Star Bullet
                        run_bullet = p.add_run()
                        run_bullet.text = "✦  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(28)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(14, 165, 233) # Sky-blue Bullet (#0EA5E9)
                    else:
                        p.space_after = Pt(8)
                        p.left_indent = Inches(0.5)
                        
                        # Dot Bullet
                        run_bullet = p.add_run()
                        run_bullet.text = "•  "
                        run_bullet.font.name = 'Microsoft JhengHei'
                        run_bullet.font.size = Pt(25)
                        run_bullet.font.bold = True
                        run_bullet.font.color.rgb = RGBColor(125, 211, 252) # Sky-blue Bullet (#7DD3FC)
                    
                    text_content = item['text']
                    if item['type'] == 'num_list':
                        text_content = f"{item['num']}. {text_content}"
                        
                    # Support **bold** formatting runs inside text
                    parts = re.split(r'(\*\*.*?\*\*)', text_content)
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            run = p.add_run()
                            run.text = part[2:-2]
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(56, 189, 248) # Sky-blue Accent for bold text (#38BDF8)
                        else:
                            run = p.add_run()
                            run.text = part
                            run.font.bold = False
                            run.font.color.rgb = RGBColor(226, 232, 240) # Slate-200 Light Gray regular text
                            
                    p.font.name = 'Microsoft JhengHei'
                    p.font.size = Pt(28 - item['level'] * 3)
                    
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

def split_slides(slides_data):
    new_slides = []
    for slide in slides_data:
        if slide['is_cover'] or not slide['items']:
            new_slides.append(slide)
            continue
        
        items = slide['items']
        chunks = []
        current_chunk = []
        current_lines = 0
        
        for item in items:
            text = item['text']
            clean_text = text.replace('**', '')
            char_limit = 43 if item['level'] > 0 else 38
            item_lines = max(1, (len(clean_text) + char_limit - 1) // char_limit)
            
            if current_lines + item_lines > 8 and current_chunk:
                chunks.append(current_chunk)
                current_chunk = [item]
                current_lines = item_lines
            else:
                current_chunk.append(item)
                current_lines += item_lines
                
        if current_chunk:
            chunks.append(current_chunk)
            
        total_pages = len(chunks)
        if total_pages <= 1:
            new_slides.append(slide)
            continue
            
        for idx, chunk in enumerate(chunks):
            page_num = idx + 1
            original_title = slide['title']
            new_title = f"{original_title} ({page_num}/{total_pages})" if original_title else f"({page_num}/{total_pages})"
            
            new_slides.append({
                'title': new_title,
                'subtitle': slide['subtitle'],
                'is_cover': slide['is_cover'],
                'items': chunk,
                'notes': slide['notes']
            })
            
    return new_slides


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python md_to_pptx.py <input_md_file> <output_pptx_file>")
        # Default fallback for this project
        input_file = r"不用版控的\艾莫 0611 鑽石營公開課_簡報.md"
        output_file = r"不用版控的\艾莫 0611 鑽石營公開課.pptx"
        print(f"Using default: {input_file} -> {output_file}")
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        
    slides_data = parse_markdown_to_slides(input_file)
    if slides_data:
        slides_data = split_slides(slides_data)
        create_presentation(slides_data, output_file)
