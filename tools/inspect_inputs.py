import sys
import io
import os
from pptx import Presentation

# Open log file for writing
log_file = open("tools/inspection_log.txt", "w", encoding="utf-8")

def log_print(msg):
    log_file.write(msg + "\n")

def inspect_pptx(path, name):
    log_print(f"\n=========================================")
    log_print(f"Inspecting: {name} ({path})")
    if not os.path.exists(path):
        log_print("File not found!")
        return
    prs = Presentation(path)
    log_print(f"Total slides: {len(prs.slides)}")
    
    for i, slide in enumerate(prs.slides):
        log_print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            shape_type = shape.shape_type
            shape_name = shape.name
            left = shape.left.inches if shape.left else 0
            top = shape.top.inches if shape.top else 0
            width = shape.width.inches if shape.width else 0
            height = shape.height.inches if shape.height else 0
            text_preview = ""
            if shape.has_text_frame:
                texts = []
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())
                text_preview = " | Texts: " + " / ".join(texts)[:200]
            log_print(f"  Shape: '{shape_name}', Type: {shape_type}, Pos: L={left:.2f}\", T={top:.2f}\", W={width:.2f}\", H={height:.2f}\"{text_preview}")

log_print("Starting PPTX inspection...")
inspect_pptx("最終輸出用檔案/圖表重製.pptx", "圖表重製")
inspect_pptx("最終輸出用檔案/陳修平短影音保姆級教學.pptx", "陳修平短影音保姆級教學")
inspect_pptx("最終輸出用檔案/行銷界九陽神功_完整講稿_排版美編完成.pptx", "行銷界九陽神功_排版美編完成")

log_file.close()
print("Inspection written to tools/inspection_log.txt successfully.")
inspect_pptx("最終輸出用檔案/圖表重製.pptx", "圖表重製")
inspect_pptx("最終輸出用檔案/陳修平短影音保姆級教學.pptx", "陳修平短影音保姆級教學")
inspect_pptx("最終輸出用檔案/行銷界九陽神功_完整講稿_排版美編完成.pptx", "行銷界九陽神功_排版美編完成")
