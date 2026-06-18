import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import io
from PIL import Image

pdf_path = r"d:\Git\true-soul-light\不用版控的\創業管理與商業模式創新.pdf"
pptx_path = r"d:\Git\true-soul-light\不用版控的\創業管理與商業模式創新.pptx"

# Initialize presentation
prs = Presentation()
# Set slide size to 10 x 7.5 inches (matches 720 x 540 pt)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # Blank layout

doc = fitz.open(pdf_path)

print(f"Starting conversion of {len(doc)} pages...")

for page_idx, page in enumerate(doc):
    print(f"Processing Page {page_idx + 1}...")
    slide = prs.slides.add_slide(blank_layout)
    
    # 1. Extract and add all images first (so they are in the background)
    image_list = page.get_images(full=True)
    for img_info in image_list:
        xref = img_info[0]
        rects = page.get_image_rects(xref)
        if not rects:
            continue
            
        try:
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            # Use PIL to convert to PNG to ensure compatibility with PPTX
            img = Image.open(io.BytesIO(image_bytes))
            img_io = io.BytesIO()
            img.save(img_io, format="PNG")
            img_io.seek(0)
            
            for r in rects:
                left = Inches(r.x0 / 72.0)
                top = Inches(r.y0 / 72.0)
                width = Inches(r.width / 72.0)
                height = Inches(r.height / 72.0)
                slide.shapes.add_picture(img_io, left, top, width, height)
        except Exception as e:
            print(f"  Warning: failed to extract/add image xref {xref} on Page {page_idx + 1}: {e}")

    # 2. Extract and add all text blocks on top
    text_page = page.get_text("dict")
    for block in text_page["blocks"]:
        if block["type"] == 0:  # Text block
            x0, y0, x1, y1 = block["bbox"]
            
            # Convert coordinates to inches
            left = Inches(x0 / 72.0)
            top = Inches(y0 / 72.0)
            # Add a small width buffer (5%) to prevent unexpected line wraps in PPTX
            width = Inches(((x1 - x0) * 1.05) / 72.0)
            height = Inches(((y1 - y0) * 1.1) / 72.0)
            
            # Avoid creating zero-size text boxes
            if width <= 0 or height <= 0:
                continue
                
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            # Remove margins for exact alignment
            tf.margin_left = Inches(0)
            tf.margin_right = Inches(0)
            tf.margin_top = Inches(0)
            tf.margin_bottom = Inches(0)
            
            for line_idx, line in enumerate(block["lines"]):
                # Create paragraph
                if line_idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Check for bullet point symbols or lists
                # (Often PDF bullets are represented in wingdings, but let's check)
                for span in line["spans"]:
                    span_text = span["text"]
                    
                    # Wingdings bullet mapping to standard bullet character
                    if "Wingdings" in span["font"]:
                        # 0xf075 or standard wingdings square/dot is often used
                        span_text = "◆ "
                    
                    run = p.add_run()
                    run.text = span_text
                    
                    # Font name mapping
                    font_name = span["font"]
                    if "DFKaiShu" in font_name or "BiauKai" in font_name or "Ming" in font_name:
                        run.font.name = "微軟正黑體"
                    else:
                        run.font.name = "微軟正黑體"  # Keep uniform Microsoft JhengHei for modern look
                        
                    # Font size
                    run.font.size = Pt(span["size"])
                    
                    # Bold & Italic
                    if "Bold" in font_name or "-Bold" in font_name or "bold" in font_name.lower():
                        run.font.bold = True
                    if "Italic" in font_name or "-Italic" in font_name or "italic" in font_name.lower():
                        run.font.italic = True
                        
                    # Color
                    color_int = span["color"]
                    r_color = (color_int >> 16) & 0xFF
                    g_color = (color_int >> 8) & 0xFF
                    b_color = color_int & 0xFF
                    run.font.color.rgb = RGBColor(r_color, g_color, b_color)

# Save the presentation
prs.save(pptx_path)
print(f"Successfully saved presentation to {pptx_path}!")
doc.close()
